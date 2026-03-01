"""
基于 output/proposal.md 生成政府汇报PPT。
项目：通州分局双中心业务运行安全监测系统

幻灯片结构（共18页）：
 1  封面
 2  汇报提纲
 3  一、项目概述（基本情况+建设目标+三子项）
 4  二、问题现状与必要性
 5  三、系统定位与边界
 6  四、总体架构与技术方案
 7  图件：系统总体架构图
 8  图件：服务器角色划分图
 9  图件：网络拓扑图
10  图件：数据流转图
11  图件：部署结构图
12  五、子项一—派出所定时打卡管理
13  六、子项二—所内视频行为分析
14  七、子项三—涉案财物视频治理算力节点
15  八、实施计划
16  图件：实施甘特图
17  九、投资预算
18  十、风险管控与保障
19  十一、建设成效与结论
20  结束页
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE      = Path(__file__).parent
OUTPUT    = BASE / "建设方案汇报.pptx"
CHART_DIR = BASE / "chart"

CHARTS = {
    "system_architecture":  CHART_DIR / "system_architecture.png",
    "server_role":          CHART_DIR / "server_role_structure.png",
    "network_topology":     CHART_DIR / "network_topology.png",
    "data_flow":            CHART_DIR / "data_flow.png",
    "deployment":           CHART_DIR / "deployment_structure.png",
    "gantt":                CHART_DIR / "project_gantt.png",
}

# ─── 色彩方案（政府汇报：深蓝+白+浅灰） ───
C_DARK   = RGBColor(0x1A, 0x2A, 0x4E)
C_ACCENT = RGBColor(0x2E, 0x5C, 0x9A)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x22, 0x22, 0x22)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_LIGHT  = RGBColor(0xEE, 0xF2, 0xF8)
C_GREEN  = RGBColor(0x19, 0x87, 0x54)
C_RED    = RGBColor(0xB5, 0x2A, 0x1D)
C_GOLD   = RGBColor(0xD4, 0x9E, 0x2A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ════════════════════════════════════════
# 基础绘图工具
# ════════════════════════════════════════

def _rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def _tb(slide, l, t, w, h, text="", size=14, color=C_TEXT,
        bold=False, align=PP_ALIGN.LEFT, wrap=True, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.name  = font
    return tf


def _para(tf, text, size=13, color=C_TEXT, bold=False,
          align=PP_ALIGN.LEFT, sb=Pt(4), sa=Pt(2)):
    from pptx.util import Pt as _Pt
    p = tf.add_paragraph()
    p.alignment   = align
    p.space_before = sb
    p.space_after  = sa
    run = p.add_run()
    run.text       = text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.name  = "Microsoft YaHei"
    return p


def _bullet(tf, text, size=13, color=C_TEXT, bold=False, prefix="● "):
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    p.space_after  = Pt(2)
    run = p.add_run()
    run.text       = prefix + text
    run.font.size  = Pt(size)
    run.font.color.rgb = color
    run.font.bold  = bold
    run.font.name  = "Microsoft YaHei"
    return p


def _title_bar(slide, title):
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), fill=C_DARK)
    _rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.045), fill=C_ACCENT)
    _tb(slide, Inches(0.55), Inches(0.14), Inches(11.5), Inches(0.72),
        text=title, size=27, color=C_WHITE, bold=True)


def _content_slide(prs, title, fn):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, title)
    fn(slide)
    return slide


def _chart_slide(prs, title, chart_key):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, title)
    path = CHARTS[chart_key]
    if path.exists():
        from PIL import Image
        with Image.open(str(path)) as img:
            iw, ih = img.size
        max_w = Inches(12.4)
        max_h = Inches(5.9)
        ratio  = min(max_w / (iw * 914400 / 96),
                     max_h / (ih * 914400 / 96))
        pw = int(iw * 914400 / 96 * ratio)
        ph = int(ih * 914400 / 96 * ratio)
        left = int((SLIDE_W - pw) / 2)
        top  = Inches(1.2)
        slide.shapes.add_picture(str(path), left, top, pw, ph)
    else:
        _tb(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.8),
            text=f"【图件缺失：{path.name}】",
            size=18, color=C_RED, align=PP_ALIGN.CENTER)
    return slide


def _card(slide, x, y, w, h, title, body_lines, title_color=C_DARK,
          bg=C_LIGHT, title_size=14, body_size=12):
    _rect(slide, x, y, w, h, fill=bg)
    _tb(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.45),
        text=title, size=title_size, color=title_color, bold=True,
        align=PP_ALIGN.CENTER)
    tf = _tb(slide, x + Inches(0.1), y + Inches(0.52), w - Inches(0.2),
             h - Inches(0.6), text="", size=body_size, color=C_GRAY,
             align=PP_ALIGN.CENTER)
    for i, line in enumerate(body_lines):
        if i == 0:
            tf.paragraphs[0].runs[0].text = line if tf.paragraphs[0].runs else None
            if tf.paragraphs[0].runs:
                tf.paragraphs[0].runs[0].text = line
        else:
            _para(tf, line, size=body_size, color=C_GRAY,
                  align=PP_ALIGN.CENTER, sb=Pt(1), sa=Pt(1))


# ════════════════════════════════════════
# 封面
# ════════════════════════════════════════

def _cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 深蓝全背景
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_DARK)
    # 顶部装饰条
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill=C_ACCENT)
    # 底部装饰条
    _rect(slide, Inches(0), SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill=C_ACCENT)
    # 中部分隔线
    _rect(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.04), fill=C_ACCENT)
    _rect(slide, Inches(1.5), Inches(5.05), Inches(10.3), Inches(0.04), fill=C_ACCENT)
    # 主标题
    _tb(slide, Inches(1.0), Inches(2.75), Inches(11.3), Inches(1.3),
        text="通州分局双中心业务运行安全监测系统",
        size=38, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    # 副标题
    _tb(slide, Inches(1.5), Inches(4.05), Inches(10.3), Inches(0.8),
        text="建 设 方 案 汇 报",
        size=28, color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
    # 底部信息
    _tb(slide, Inches(1.5), Inches(5.25), Inches(10.3), Inches(0.5),
        text="建设单位：北京市公安局通州分局　　　建设周期：6个月　　　总投资：91.7735万元",
        size=14, color=RGBColor(0x88, 0x99, 0xBB), align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.4),
        text="文件状态：交付稿　　密级等级：内部　　版本号：V1.0",
        size=13, color=RGBColor(0x66, 0x77, 0x99), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 汇报提纲
# ════════════════════════════════════════

def _toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _title_bar(slide, "汇报提纲")
    items_l = [
        "一、项目概述",
        "二、问题现状与必要性",
        "三、系统定位与边界",
        "四、总体架构与技术方案",
        "五、子项一—派出所定时打卡管理",
    ]
    items_r = [
        "六、子项二—所内视频行为分析",
        "七、子项三—涉案财物视频治理算力节点",
        "八、实施计划",
        "九、投资预算",
        "十、风险管控、保障与建设成效",
    ]
    for col, items in enumerate([items_l, items_r]):
        x = Inches(0.6 + col * 6.5)
        for row, item in enumerate(items):
            y = Inches(1.4 + row * 1.05)
            _rect(slide, x, y, Inches(5.8), Inches(0.85), fill=C_LIGHT)
            _rect(slide, x, y, Inches(0.08), Inches(0.85), fill=C_ACCENT)
            _tb(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.4), Inches(0.55),
                text=item, size=18, color=C_DARK, bold=True)


# ════════════════════════════════════════
# 第1页：项目概述
# ════════════════════════════════════════

def _p1(slide):
    # 基本情况框
    _rect(slide, Inches(0.4), Inches(1.15), Inches(5.8), Inches(3.1), fill=C_LIGHT)
    _tb(slide, Inches(0.55), Inches(1.2), Inches(3), Inches(0.45),
        text="项目基本情况", size=17, color=C_ACCENT, bold=True)
    tf = _tb(slide, Inches(0.55), Inches(1.65), Inches(5.5), Inches(2.4),
             text="", size=13, color=C_TEXT)
    items = [
        "项目名称：通州分局双中心业务运行安全监测系统",
        "建设单位：北京市公安局通州分局",
        "建设范围：执法办案中心、涉案财物管理中心及辖区27个派出所",
        "项目总投资：91.7735万元（人民币）",
        "建设周期：6个月　资金来源：财政资金",
    ]
    for i, t in enumerate(items):
        if i == 0:
            tf.paragraphs[0].runs[0].text = t if tf.paragraphs[0].runs else t
            if not tf.paragraphs[0].runs:
                _para(tf, t, size=13, sb=Pt(3), sa=Pt(2))
        else:
            _para(tf, t, size=13, sb=Pt(3), sa=Pt(2))

    # 建设目标框
    _rect(slide, Inches(6.5), Inches(1.15), Inches(6.4), Inches(3.1), fill=C_LIGHT)
    _tb(slide, Inches(6.65), Inches(1.2), Inches(3.5), Inches(0.45),
        text="建设目标", size=17, color=C_ACCENT, bold=True)
    tf2 = _tb(slide, Inches(6.65), Inches(1.65), Inches(6.1), Inches(2.4),
              text="", size=13, color=C_TEXT)
    goals = [
        "看管岗位履职状态持续可见",
        "涉案财物管理全过程留痕与追溯",
        "关键行为事件结构化记录与管理核查支撑",
        '视频资源从"存储资料"向"管理能力载体"功能转化',
        '不追求"自动化执法"，所有事件须经人工核查',
    ]
    for i, g in enumerate(goals):
        if i == 0:
            if tf2.paragraphs[0].runs:
                tf2.paragraphs[0].runs[0].text = "● " + g
        else:
            _bullet(tf2, g, size=13)

    # 三子项简介
    _rect(slide, Inches(0.4), Inches(4.45), SLIDE_W - Inches(0.8), Inches(0.04),
          fill=C_ACCENT)
    _tb(slide, Inches(0.4), Inches(4.6), Inches(5), Inches(0.4),
        text="三个建设子项", size=17, color=C_ACCENT, bold=True)
    sub_items = [
        ("子项一", "派出所定时打卡管理", "27个派出所部署人脸识别终端，统一打卡管理", "23.51万"),
        ("子项二", "所内视频行为分析",   "54台岗位摄像头+AI分析，看管履职疑似行为检测", "26.26万"),
        ("子项三", "涉案财物视频治理",   "1台GPU服务器（DCU K100-AI×2），承载28路视频分析", "42.00万"),
    ]
    for i, (tag, title, desc, amt) in enumerate(sub_items):
        x = Inches(0.4 + i * 4.3)
        _rect(slide, x, Inches(5.1), Inches(4.0), Inches(2.1), fill=C_LIGHT)
        _rect(slide, x, Inches(5.1), Inches(0.08), Inches(2.1), fill=C_GOLD)
        _tb(slide, x + Inches(0.18), Inches(5.15), Inches(3.7), Inches(0.4),
            text=f"{tag}  {title}", size=14, color=C_DARK, bold=True)
        _tb(slide, x + Inches(0.18), Inches(5.55), Inches(3.7), Inches(1.2),
            text=desc, size=12, color=C_GRAY, wrap=True)
        _tb(slide, x + Inches(0.18), Inches(6.7), Inches(3.7), Inches(0.35),
            text=f"预算：{amt}", size=12, color=C_ACCENT, bold=True)


# ════════════════════════════════════════
# 第2页：问题现状与必要性
# ════════════════════════════════════════

def _p2(slide):
    # 三大问题
    _tb(slide, Inches(0.4), Inches(1.15), Inches(6), Inches(0.45),
        text="当前管理三大核心问题", size=17, color=C_ACCENT, bold=True)
    problems = [
        ("履职监督空白",
         ["缺乏持续性技术监测手段", "脱岗等疑似行为无法及时发现", "人工监督时效性和连续性有限"]),
        ("财物管理核查低效",
         ["人员活动依赖人工翻查视频", "双人操作规则核查成本高", "无法形成结构化管理台账"]),
        ("管理规则落地困难",
         ["双人操作、非工作时间入室管控", "无技术手段支撑，落地效果差", "27个派出所一致性难保障"]),
    ]
    for i, (t, lines) in enumerate(problems):
        x = Inches(0.4 + i * 4.3)
        _rect(slide, x, Inches(1.65), Inches(4.0), Inches(2.6), fill=C_LIGHT)
        _rect(slide, x, Inches(1.65), Inches(4.0), Inches(0.5),
              fill=RGBColor(0xB5, 0x2A, 0x1D))
        _tb(slide, x + Inches(0.15), Inches(1.7), Inches(3.7), Inches(0.42),
            text=t, size=15, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        tf = _tb(slide, x + Inches(0.15), Inches(2.2), Inches(3.7), Inches(1.9),
                 text="", size=13)
        for j, ln in enumerate(lines):
            if j == 0:
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = "● " + ln
            else:
                _bullet(tf, ln, size=13)

    # 分隔线
    _rect(slide, Inches(0.4), Inches(4.45), SLIDE_W - Inches(0.8),
          Inches(0.04), fill=C_ACCENT)
    _tb(slide, Inches(0.4), Inches(4.6), Inches(5), Inches(0.45),
        text="建设必要性（7条）", size=17, color=C_ACCENT, bold=True)

    necessity = [
        "履职监管需要：建立持续性、客观性履职状态技术监测能力",
        "财物合规需要：形成可核查的结构化记录，支撑审计与核查",
        "视频效用提升：将视频从存储资料转化为主动管理工具",
        "管理统一性：统一规则框架，保障27个派出所管理逻辑一致性",
        "风险提前发现：异常行为事前提示，主动防控替代被动应对",
        "信息化建设：填补通州分局在看管安全监管方向的建设空白",
        "减少核查成本：释放人力，提升核查效率与管理覆盖率",
    ]
    # 分两列显示
    for i, n in enumerate(necessity):
        col = i % 2
        row = i // 2
        x = Inches(0.4 + col * 6.5)
        y = Inches(5.15 + row * 0.52)
        _tb(slide, x, y, Inches(6.2), Inches(0.48),
            text=f"{'①②③④⑤⑥⑦'[i]}  {n}", size=13, color=C_TEXT)


# ════════════════════════════════════════
# 第3页：系统定位与边界
# ════════════════════════════════════════

def _p3(slide):
    # 定位说明
    _rect(slide, Inches(0.4), Inches(1.15), SLIDE_W - Inches(0.8), Inches(1.3),
          fill=C_LIGHT)
    _rect(slide, Inches(0.4), Inches(1.15), Inches(0.1), Inches(1.3), fill=C_ACCENT)
    _tb(slide, Inches(0.65), Inches(1.22), Inches(11.8), Inches(0.45),
        text="系统定位：辅助管理能力建设，而非执法能力建设",
        size=18, color=C_DARK, bold=True)
    _tb(slide, Inches(0.65), Inches(1.68), Inches(11.8), Inches(0.65),
        text="工具定位：内部管理辅助工具，面向分局管理层及派出所负责人，产生的事件记录供管理人员内部核查使用，不作为直接处罚凭据。",
        size=13, color=C_TEXT, wrap=True)

    # 建设边界表
    _tb(slide, Inches(0.4), Inches(2.65), Inches(4), Inches(0.45),
        text="建设边界", size=17, color=C_ACCENT, bold=True)
    headers = ["边界维度", "本次建设范围内", "明确排除"]
    rows = [
        ("系统类型",   "监测辅助类系统",             "业务主系统、执法决定系统"),
        ("数据类型",   "事件级短视频（10–30秒）",     "原始全量视频流集中存储"),
        ("结果性质",   "管理辅助信息，须人工确认",    "自动化违规认定、直接处罚依据"),
        ("既有系统",   "接入执法办案2.0信号（只读）",  "修改既有业务系统任何逻辑"),
        ("网络范围",   "公安网内",                   "互联网、区政务云"),
        ("识别技术",   "打卡方向：人脸识别仅用于打卡", "看管检测方向不做人脸识别"),
    ]
    col_widths = [Inches(2.2), Inches(4.5), Inches(5.7)]
    # 表头
    header_fill = RGBColor(0x2E, 0x5C, 0x9A)
    for j, (h, cw) in enumerate(zip(headers, col_widths)):
        x = Inches(0.4) + sum(col_widths[:j])
        _rect(slide, x, Inches(3.15), cw, Inches(0.45), fill=header_fill,
              line=C_WHITE, line_w=Pt(1))
        _tb(slide, x + Inches(0.1), Inches(3.18), cw - Inches(0.2), Inches(0.38),
            text=h, size=13, color=C_WHITE, bold=True)
    # 数据行
    row_colors = [C_LIGHT, C_WHITE]
    for i, (d, v1, v2) in enumerate(rows):
        y = Inches(3.6 + i * 0.52)
        bg = row_colors[i % 2]
        for j, (val, cw) in enumerate(zip([d, v1, v2], col_widths)):
            x = Inches(0.4) + sum(col_widths[:j])
            _rect(slide, x, y, cw, Inches(0.5), fill=bg,
                  line=RGBColor(0xCC, 0xD5, 0xE5), line_w=Pt(1))
            _tb(slide, x + Inches(0.1), y + Inches(0.07), cw - Inches(0.2),
                Inches(0.38), text=val, size=12,
                color=C_DARK if j == 0 else C_TEXT)


# ════════════════════════════════════════
# 第4页：总体架构与技术方案
# ════════════════════════════════════════

def _p4(slide):
    # 四层架构
    _tb(slide, Inches(0.4), Inches(1.15), Inches(5), Inches(0.45),
        text="系统四层架构", size=17, color=C_ACCENT, bold=True)
    layers = [
        ("采集层", "人脸识别终端（27台）、岗位值守摄像头（54台）、保管室摄像头",
         RGBColor(0x5B, 0x8D, 0xCE)),
        ("处理层", "服务器-01（打卡管理）、服务器-02（行为分析）、服务器-03（GPU视频分析）",
         RGBColor(0x2E, 0x5C, 0x9A)),
        ("存储层", "各服务器本地HDD，RAID等价保护，事件级短视频+结构化记录",
         RGBColor(0x19, 0x87, 0x54)),
        ("管理层", "统一管理视图—告警管理、统计分析、记录查询、权限配置，公安网内访问",
         C_DARK),
    ]
    for i, (name, desc, color) in enumerate(layers):
        y = Inches(1.65 + i * 1.15)
        _rect(slide, Inches(0.4), y, SLIDE_W * 0.52, Inches(1.05), fill=C_LIGHT)
        _rect(slide, Inches(0.4), y, Inches(1.5), Inches(1.05), fill=color)
        _tb(slide, Inches(0.5), y + Inches(0.28), Inches(1.3), Inches(0.5),
            text=name, size=15, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, Inches(2.1), y + Inches(0.12), Inches(4.7), Inches(0.82),
            text=desc, size=12, color=C_TEXT, wrap=True)

    # 服务器角色对比（右栏）
    _tb(slide, Inches(7.3), Inches(1.15), Inches(5.5), Inches(0.45),
        text="三台服务器角色", size=17, color=C_ACCENT, bold=True)
    servers = [
        ("服务器-01", "打卡管理",
         ["海光 Hygon 7360×2（24C）", "128GB DDR4，1+2×8TB HDD", "配套打卡管理软件"]),
        ("服务器-02", "视频行为分析",
         ["海光 Hygon 7360×2（24C）", "128GB DDR4，1+2×8TB HDD", "配套行为分析软件（5万元）"]),
        ("服务器-03", "GPU视频分析",
         ["海光 Hygon 7470×2（48C）", "512GB DDR5，DCU K100-AI×2", "承载28路并发视频分析"]),
    ]
    for i, (sn, role, specs) in enumerate(servers):
        y = Inches(1.65 + i * 1.85)
        _rect(slide, Inches(7.3), y, Inches(5.6), Inches(1.7), fill=C_LIGHT)
        _rect(slide, Inches(7.3), y, Inches(1.4), Inches(1.7),
              fill=C_DARK if i == 2 else C_ACCENT)
        _tb(slide, Inches(7.35), y + Inches(0.15), Inches(1.3), Inches(0.5),
            text=sn, size=12, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
        _tb(slide, Inches(7.35), y + Inches(0.65), Inches(1.3), Inches(0.5),
            text=role, size=11, color=RGBColor(0xCC, 0xDD, 0xFF),
            align=PP_ALIGN.CENTER)
        tf = _tb(slide, Inches(8.85), y + Inches(0.1), Inches(3.9), Inches(1.5),
                 text="", size=12)
        for j, sp in enumerate(specs):
            if j == 0:
                if tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = sp
            else:
                _para(tf, sp, size=12, sb=Pt(2), sa=Pt(1))

    _tb(slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.45),
        text="（详见下页图件：系统总体架构图、服务器角色划分图、网络拓扑图、数据流转图、部署结构图）",
        size=13, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 第5页：子项一—定时打卡管理
# ════════════════════════════════════════

def _p5(slide):
    _tb(slide, Inches(0.4), Inches(1.15), Inches(9), Inches(0.45),
        text="建设内容", size=17, color=C_ACCENT, bold=True)
    tf = _tb(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.6),
             text="", size=13)
    content = [
        "在27个派出所各部署1套人脸识别终端（7英寸触摸屏、双目摄像头，支持WiFi/TCP/IP）及配套设施（报警器、延时继电器、电源适配器、线材线缆）",
        "部署1台打卡管理应用服务器（海光 Hygon 7360×2，128GB DDR4），配套打卡管理软件（10,000元）",
        "工程路径：基于刷脸门禁设备，通过时间继电器与信号联动形成定点巡更打卡闭环（已验证工程路径，规模化复用）",
    ]
    for i, c in enumerate(content):
        if i == 0:
            if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = "● " + c
        else:
            _bullet(tf, c, size=13)

    _rect(slide, Inches(0.4), Inches(3.35), SLIDE_W - Inches(0.8),
          Inches(0.04), fill=C_ACCENT)
    _tb(slide, Inches(0.4), Inches(3.5), Inches(5), Inches(0.45),
        text="核心功能", size=17, color=C_ACCENT, bold=True)
    funcs = [
        ("角色分级打卡", "按岗位角色设置不同打卡时间间隔，满足不同岗位履职要求"),
        ("未打卡提醒",   "声光报警联动，超时未打卡自动触发报警，通知机制可配置"),
        ("打卡记录管理", "按日期、派出所、岗位查询，支持记录导出，完整留存"),
        ("统计分析视图", "打卡达标率、异常记录汇总报表，分局层面汇总视图"),
        ("系统管理功能", "设备管理（状态查看、远程配置）、用户角色管理、参数配置"),
    ]
    for i, (fn, desc) in enumerate(funcs):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(4.05 + row * 1.65)
        _rect(slide, x, y, Inches(4.0), Inches(1.5), fill=C_LIGHT)
        _rect(slide, x, y, Inches(0.08), Inches(1.5), fill=C_GOLD)
        _tb(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.45),
            text=fn, size=14, color=C_DARK, bold=True)
        _tb(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(0.85),
            text=desc, size=12, color=C_GRAY, wrap=True)

    # 预算
    _rect(slide, Inches(12.0), Inches(3.5), Inches(1.0), Inches(0.5),
          fill=C_ACCENT)
    _tb(slide, Inches(11.0), Inches(3.5), Inches(2.2), Inches(0.48),
        text="预算：235,145元", size=14, color=C_ACCENT, bold=True,
        align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════
# 第6页：子项二—视频行为分析
# ════════════════════════════════════════

def _p6(slide):
    _tb(slide, Inches(0.4), Inches(1.15), Inches(9), Inches(0.45),
        text="建设内容", size=17, color=C_ACCENT, bold=True)
    tf = _tb(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.6),
             text="", size=13)
    content = [
        "在27个派出所执法办案中心等候区各部署2台岗位值守专用摄像头（DS-2XA8D45F/LG-IZS，400万像素，内置AI分析），共54台",
        "部署1台视频行为分析应用服务器（海光 Hygon 7360×2，128GB DDR4），配套行为分析软件（50,000元）",
        '架构：摄像头内置AI初步分析，服务器侧事件汇聚+逻辑判断+管理承载，"端+云"轻量化处理',
    ]
    for i, c in enumerate(content):
        if i == 0:
            if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = "● " + c
        else:
            _bullet(tf, c, size=13)

    _rect(slide, Inches(0.4), Inches(3.35), SLIDE_W - Inches(0.8),
          Inches(0.04), fill=C_ACCENT)
    _tb(slide, Inches(0.4), Inches(3.5), Inches(5), Inches(0.45),
        text="核心功能", size=17, color=C_ACCENT, bold=True)
    funcs = [
        ("嫌疑人在场判断",
         "双模式：执法办案2.0信号（精确）+ 视频检测+时间阈值（≥5分钟入区），仅在嫌疑人在场时激活检测"),
        ("疑似行为检测",
         "固定ROI区域内连续分析脱岗（连续无人≥5分钟）、睡岗、玩手机，不做人脸识别"),
        ("告警与证据留存",
         "告警触发≤30秒推送，关联时间戳+摄像头+关键图片+10–30秒事件视频"),
        ("人工确认机制",
         "支持确认、驳回、标注操作，每条告警关联操作记录，检测结果定性为辅助管理信息"),
        ("统计分析视图",
         "按时段、派出所、事件类型汇总，支持分局层面横向比较与周期性报表"),
    ]
    for i, (fn, desc) in enumerate(funcs):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(4.05 + row * 1.65)
        _rect(slide, x, y, Inches(4.0), Inches(1.5), fill=C_LIGHT)
        _rect(slide, x, y, Inches(0.08), Inches(1.5), fill=C_ACCENT)
        _tb(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.45),
            text=fn, size=14, color=C_DARK, bold=True)
        _tb(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(0.85),
            text=desc, size=12, color=C_GRAY, wrap=True)

    _tb(slide, Inches(11.0), Inches(3.5), Inches(2.2), Inches(0.48),
        text="预算：262,590元", size=14, color=C_ACCENT, bold=True,
        align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════
# 第7页：子项三—涉案财物视频治理
# ════════════════════════════════════════

def _p7(slide):
    _tb(slide, Inches(0.4), Inches(1.15), Inches(9), Inches(0.45),
        text="建设内容", size=17, color=C_ACCENT, bold=True)
    tf = _tb(slide, Inches(0.4), Inches(1.65), Inches(12.5), Inches(1.6),
             text="", size=13)
    content = [
        "部署1台视频分析服务器（海光 Hygon 7470×2/48C，512GB ECC DDR5，DCU K100-AI 64GB×2，2×960GB SSD+2×8TB HDD，1+1冗余电源），含配套视频分析功能软件",
        "承载分局涉案财物管理中心（1处）及27个派出所保管室（27处）的视频AI分析，合计28路并发分析",
        "预算：420,000元（含硬件+软件集成交付，最大单项投资）",
    ]
    for i, c in enumerate(content):
        if i == 0:
            if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = "● " + c
        else:
            _bullet(tf, c, size=13)

    _rect(slide, Inches(0.4), Inches(3.35), SLIDE_W - Inches(0.8),
          Inches(0.04), fill=C_ACCENT)
    _tb(slide, Inches(0.4), Inches(3.5), Inches(5), Inches(0.45),
        text="核心功能", size=17, color=C_ACCENT, bold=True)
    funcs = [
        ("人员活动留痕",
         "出现时间+停留时长+区域关联+关键帧+10–30秒事件短视频，全时段持续分析"),
        ("非工作时间标记",
         "默认工作时间09:00–18:00，非工作时间内活动自动加注标记，参数可配置"),
        ("出入库人数校验",
         "区域人数辅助校验，少于2人时触发提示，辅助双人操作规则落地"),
        ("异常行为提示",
         "异常停留超阈值、短时重复进出、偏离历史模式等情形自动提示"),
        ("视频辅助台账视图",
         "时间轴呈现事件序列，与业务台账时间节点关联查阅"),
        ("查询导出",
         "按时间、区域、事件类型查询，支持结构化记录与关键帧导出，满足审计核查"),
    ]
    for i, (fn, desc) in enumerate(funcs):
        col = i % 3
        row = i // 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(4.05 + row * 1.65)
        _rect(slide, x, y, Inches(4.0), Inches(1.5), fill=C_LIGHT)
        _rect(slide, x, y, Inches(0.08), Inches(1.5), fill=C_GREEN)
        _tb(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.45),
            text=fn, size=14, color=C_DARK, bold=True)
        _tb(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(0.85),
            text=desc, size=12, color=C_GRAY, wrap=True)


# ════════════════════════════════════════
# 第8页：实施计划
# ════════════════════════════════════════

def _p8(slide):
    _tb(slide, Inches(0.4), Inches(1.15), Inches(9), Inches(0.45),
        text="三阶段实施计划（建设周期：6个月）", size=17, color=C_ACCENT, bold=True)
    phases = [
        ("第一阶段\n（1–2月）",
         "设备采购与基础部署",
         ["服务器、摄像头、人脸识别终端采购到货",
          "机房上架，基础环境配置",
          "操作系统、数据库基础部署",
          "国产化环境适配验证"],
         "M1：设备到货验收完成，基础环境就绪"),
        ("第二阶段\n（3–4月）",
         "软件部署与设备接入",
         ["应用软件安装部署",
          "27个派出所设备分批接入调试",
          "双模式嫌疑人在场判断调试",
          "各子系统功能测试（含GPU28路并发验证）"],
         "M2：27所设备接入完成，各子系统功能测试通过"),
        ("第三阶段\n（5–6月）",
         "联调测试与上线验收",
         ["全系统联调",
          "信息安全等保二级测评",
          "第三方软件测评（功能/性能/安全）",
          "用户培训+上线验收"],
         "M3：系统上线运行并通过验收"),
    ]
    for i, (period, title, tasks, milestone) in enumerate(phases):
        x = Inches(0.4 + i * 4.3)
        # 阶段卡片
        _rect(slide, x, Inches(1.65), Inches(4.0), Inches(4.5), fill=C_LIGHT)
        _rect(slide, x, Inches(1.65), Inches(4.0), Inches(0.75),
              fill=C_DARK if i == 2 else (C_ACCENT if i == 1 else RGBColor(0x3A, 0x70, 0xB5)))
        _tb(slide, x + Inches(0.15), Inches(1.7), Inches(3.7), Inches(0.65),
            text=f"{period}  {title}", size=14, color=C_WHITE, bold=True)
        tf = _tb(slide, x + Inches(0.15), Inches(2.5), Inches(3.7), Inches(3.3),
                 text="", size=13)
        for j, t in enumerate(tasks):
            if j == 0:
                if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = "● " + t
            else:
                _bullet(tf, t, size=13)

        # 里程碑
        _rect(slide, x, Inches(6.25), Inches(4.0), Inches(0.75),
              fill=RGBColor(0xD4, 0x9E, 0x2A))
        _tb(slide, x + Inches(0.1), Inches(6.32), Inches(3.8), Inches(0.6),
            text=milestone, size=12, color=C_WHITE, bold=True, wrap=True)

    _tb(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
        text="（详见下页：实施甘特图）", size=13, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 第9页：投资预算
# ════════════════════════════════════════

def _p9(slide):
    # 总额展示
    _rect(slide, Inches(4.2), Inches(1.15), Inches(4.9), Inches(1.15), fill=C_DARK)
    _tb(slide, Inches(4.2), Inches(1.22), Inches(4.9), Inches(0.45),
        text="项目投资总额", size=16, color=RGBColor(0xAA, 0xBB, 0xDD),
        align=PP_ALIGN.CENTER)
    _tb(slide, Inches(4.2), Inches(1.62), Inches(4.9), Inches(0.6),
        text="91.7735万元", size=34, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 三子项卡片
    cats = [
        ("子项一\n派出所定时打卡管理", "235,145元\n（23.51万元）",
         ["人脸识别终端×27  76,140元",
          "配套设施（报警器/继电器/线缆等）16,605元",
          "安装调试×27项  32,400元",
          "配套软件  10,000元",
          "打卡管理服务器×1  100,000元"],
         C_GOLD),
        ("子项二\n所内视频行为分析", "262,590元\n（26.26万元）",
         ["岗位值守摄像头×54  72,900元",
          "线材线缆×27套  7,290元",
          "安装调试×27项  32,400元",
          "行为分析软件  50,000元",
          "视频行为分析服务器×1  100,000元"],
         C_ACCENT),
        ("子项三\n涉案财物视频治理算力节点", "420,000元\n（42.00万元）",
         ["视频分析服务器（含GPU DCU K100-AI×2）×1  420,000元",
          "（含配套涉案财物安防视频分析功能）",
          "承载28路并发视频分析（分局+27派出所）",
          "最大单项投资，含全部硬件+软件集成"],
         C_GREEN),
    ]
    for i, (title, amount, details, color) in enumerate(cats):
        x = Inches(0.4 + i * 4.3)
        _rect(slide, x, Inches(2.55), Inches(4.0), Inches(4.55), fill=C_LIGHT)
        _rect(slide, x, Inches(2.55), Inches(4.0), Inches(0.75), fill=color)
        _tb(slide, x + Inches(0.15), Inches(2.6), Inches(3.7), Inches(0.65),
            text=title, size=13, color=C_WHITE, bold=True)
        _tb(slide, x + Inches(0.15), Inches(3.37), Inches(3.7), Inches(0.55),
            text=amount, size=18, color=color, bold=True)
        tf = _tb(slide, x + Inches(0.15), Inches(3.98), Inches(3.7), Inches(2.95),
                 text="", size=12)
        for j, d in enumerate(details):
            if j == 0:
                if tf.paragraphs[0].runs: tf.paragraphs[0].runs[0].text = d
            else:
                _para(tf, d, size=12, sb=Pt(3), sa=Pt(1))

    _tb(slide, Inches(0.4), Inches(7.1), Inches(12.5), Inches(0.35),
        text="资金来源：财政资金　　软件费合计：60,000元（含独立计价部分）",
        size=13, color=C_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 第10页：风险管控与保障
# ════════════════════════════════════════

def _p10(slide):
    _tb(slide, Inches(0.4), Inches(1.15), Inches(5), Inches(0.45),
        text="风险管控", size=17, color=C_ACCENT, bold=True)
    risks = [
        ("设备采购风险",    "供应商交货延迟",
         "要求到货承诺+缓冲周期；优先选择现货供应商"),
        ("技术接入风险",    "派出所网络条件差异影响调试",
         "提前网络预调查；分3–4批次接入，每批独立验证"),
        ("兼容性风险",      "国产OS/数据库与软件兼容性问题",
         "投标前提供国产环境适配说明；预留兼容性测试时间"),
        ("业务连续性风险",  "新建系统部署影响既有业务系统",
         "独立部署，不修改既有系统；调试在业务低峰期进行"),
        ("等保测评风险",    "等保测评未达二级要求",
         "部署阶段同步安全加固；整改周期预留在第三阶段"),
        ("用户接受度风险",  "基层民警使用不熟练",
         "分批培训+现场指导；配套操作手册（图文并茂）"),
    ]
    for i, (title, risk, measure) in enumerate(risks):
        row = i // 3
        col = i % 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(1.65 + row * 2.6)
        _rect(slide, x, y, Inches(4.0), Inches(2.4), fill=C_LIGHT)
        _rect(slide, x, y, Inches(4.0), Inches(0.5),
              fill=C_DARK if row == 1 else C_ACCENT)
        _tb(slide, x + Inches(0.15), y + Inches(0.06), Inches(3.7), Inches(0.4),
            text=title, size=15, color=C_WHITE, bold=True)
        _tb(slide, x + Inches(0.15), y + Inches(0.58), Inches(3.7), Inches(0.65),
            text=f"风险：{risk}", size=12, color=C_RED)
        _tb(slide, x + Inches(0.15), y + Inches(1.25), Inches(3.7), Inches(1.0),
            text=f"应对：{measure}", size=12, color=C_GREEN, wrap=True)

    # 保障机制简要
    _rect(slide, Inches(0.4), Inches(6.9), SLIDE_W - Inches(0.8),
          Inches(0.04), fill=C_ACCENT)
    _tb(slide, Inches(0.4), Inches(7.0), Inches(12.5), Inches(0.4),
        text="保障机制：第三方监理全程监理 | 等保二级测评 | 第三方软件测评 | 分批接入独立验证 | 三类角色培训",
        size=13, color=C_TEXT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 第11页：建设成效与结论
# ════════════════════════════════════════

def _p11(slide):
    _tb(slide, Inches(0.4), Inches(1.15), Inches(5), Inches(0.45),
        text="预期建设成效", size=17, color=C_ACCENT, bold=True)
    achievements = [
        ("履职监督能力",    "看管岗位疑似违规事件自动检测率 ≥90%\n（基于摄像头覆盖区域）"),
        ("财物留痕能力",    "保管室人员活动结构化记录完整率 100%\n全时段持续分析"),
        ("规则校验能力",    "非工作时间活动自动标记率 100%\n双人操作规则辅助校验覆盖率 100%"),
        ("打卡管理能力",    "打卡记录可查询完整率 100%\n含记录导出能力"),
        ("系统可用性",      "月度可用性 ≥99.5%\n告警触发至展示 ≤30秒"),
        ("查询响应能力",    "历史记录检索响应 ≤5秒\n打卡确认响应 ≤3秒"),
    ]
    for i, (title, desc) in enumerate(achievements):
        row = i // 3
        col = i % 3
        x = Inches(0.4 + col * 4.3)
        y = Inches(1.65 + row * 2.0)
        _rect(slide, x, y, Inches(4.0), Inches(1.75), fill=C_LIGHT)
        _rect(slide, x, y, Inches(0.08), Inches(1.75), fill=C_GREEN)
        _tb(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.6), Inches(0.45),
            text=title, size=15, color=C_DARK, bold=True)
        _tb(slide, x + Inches(0.2), y + Inches(0.55), Inches(3.6), Inches(1.1),
            text=desc, size=13, color=C_GRAY, wrap=True)

    _rect(slide, Inches(0.4), Inches(5.8), SLIDE_W - Inches(0.8),
          Inches(0.04), fill=C_ACCENT)
    _rect(slide, Inches(0.4), Inches(5.95), SLIDE_W - Inches(0.8),
          Inches(1.35), fill=C_LIGHT)
    _rect(slide, Inches(0.4), Inches(5.95), Inches(0.1), Inches(1.35), fill=C_DARK)
    _tb(slide, Inches(0.65), Inches(6.05), Inches(12.1), Inches(1.15),
        text=(
            "本项目聚焦执法办案与涉案财物管理两条主线，"
            "建设范围明确、目标清晰、技术路线服从 technical_baseline.md 基线、风险可控。"
            "三个子项功能独立、管理汇聚，共同构成覆盖通州分局全辖区的看管安全监测能力体系。"
            "系统定位为辅助管理工具，不追求自动化执法，所有结果须经人工核查，"
            "不改变既有业务系统，不集中原始视频数据，建设边界清晰。"
        ),
        size=14, color=C_TEXT, wrap=True)


# ════════════════════════════════════════
# 结束页
# ════════════════════════════════════════

def _end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=C_DARK)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.18), fill=C_ACCENT)
    _rect(slide, Inches(0), SLIDE_H - Inches(0.18), SLIDE_W, Inches(0.18), fill=C_ACCENT)
    _rect(slide, Inches(1.5), Inches(2.7), Inches(10.3), Inches(0.04), fill=C_ACCENT)
    _rect(slide, Inches(1.5), Inches(4.85), Inches(10.3), Inches(0.04), fill=C_ACCENT)
    _tb(slide, Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.2),
        text="汇 报 完 毕", size=48, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1.5), Inches(4.1), Inches(10.3), Inches(0.75),
        text="敬 请 指 导", size=30,
        color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
    _tb(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.5),
        text="北京市公安局通州分局",
        size=16, color=RGBColor(0x77, 0x88, 0xAA), align=PP_ALIGN.CENTER)


# ════════════════════════════════════════
# 主构建流程
# ════════════════════════════════════════

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _cover(prs)                                                    # 1  封面
    _toc(prs)                                                      # 2  汇报提纲
    _content_slide(prs, "一、项目概述",            _p1)           # 3
    _content_slide(prs, "二、问题现状与必要性",    _p2)           # 4
    _content_slide(prs, "三、系统定位与边界",      _p3)           # 5
    _content_slide(prs, "四、总体架构与技术方案",  _p4)           # 6
    _chart_slide(prs,   "系统总体架构图",           "system_architecture")   # 7
    _chart_slide(prs,   "服务器角色划分图",         "server_role")           # 8
    _chart_slide(prs,   "网络拓扑图",              "network_topology")       # 9
    _chart_slide(prs,   "数据流转图",              "data_flow")              # 10
    _chart_slide(prs,   "部署结构图",              "deployment")             # 11
    _content_slide(prs, "五、子项一—派出所定时打卡管理",          _p5)  # 12
    _content_slide(prs, "六、子项二—所内视频行为分析",            _p6)  # 13
    _content_slide(prs, "七、子项三—涉案财物视频治理算力节点",    _p7)  # 14
    _content_slide(prs, "八、实施计划",            _p8)           # 15
    _chart_slide(prs,   "实施甘特图",              "gantt")                  # 16
    _content_slide(prs, "九、投资预算",            _p9)           # 17
    _content_slide(prs, "十、风险管控与保障措施",  _p10)          # 18
    _content_slide(prs, "十一、建设成效与结论",    _p11)          # 19
    _end(prs)                                                      # 20  结束页

    prs.save(str(OUTPUT))
    print(f"Done: {OUTPUT}  （共{len(prs.slides)}页）")


if __name__ == '__main__':
    build()
