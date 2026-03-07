"""
基于 output/建设方案.md 生成政府汇报PPT。

页面结构强制约束（9页 + 图件页）：
1. 项目背景与建设必要性
2. 建设目标与原则
3. 系统定位与边界
4. 总体架构（插入架构图）
5. 数据与存储策略（插入数据流转图）
6. 部署结构（插入网络拓扑图 + 部署结构图）
7. 投资与预算
8. 风险与控制
9. 建设成效与结论

图件规则：
- 每张图件独占一页
- 图件页带标题，不加解释性长文
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = Path(__file__).parent
OUTPUT = BASE / "执法音视频集中存储与管理系统建设方案.pptx"

CHARTS = {
    "系统总体架构图": BASE / "chart" / "系统总体架构图.png",
    "网络拓扑图":     BASE / "chart" / "网络拓扑图.png",
    "数据流转图":     BASE / "chart" / "数据流转图.png",
    "部署结构图":     BASE / "chart" / "部署结构图.png",
}

# ── 颜色方案（政府汇报风格：深蓝+白+浅灰） ──
C_DARK   = RGBColor(0x1A, 0x2A, 0x4E)  # 深蓝（标题栏）
C_ACCENT = RGBColor(0x2E, 0x5C, 0x9A)  # 中蓝（装饰线/强调）
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT   = RGBColor(0x33, 0x33, 0x33)
C_GRAY   = RGBColor(0x66, 0x66, 0x66)
C_LIGHT  = RGBColor(0xF0, 0xF3, 0xF7)  # 浅灰蓝背景
C_RED    = RGBColor(0xC0, 0x39, 0x2B)

# ── 尺寸常量（16:9） ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


# ── 基础绘图工具 ──

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.color.rgb = line_color
    return shape


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_color=C_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=14, font_color=C_TEXT, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(2),
                  font_name="Microsoft YaHei"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_bullet(tf, text, font_size=14, font_color=C_TEXT, level=0, bold=False):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.level = level
    p.space_before = Pt(3)
    p.space_after = Pt(2)
    return p


# ── 页面模板 ──

def make_title_bar(slide, title_text):
    """顶部深蓝标题栏"""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.0), fill_color=C_DARK)
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(11), Inches(0.7),
                text=title_text, font_size=28, font_color=C_WHITE, bold=True,
                alignment=PP_ALIGN.LEFT)
    # 底部装饰线
    add_rect(slide, Inches(0), Inches(1.0), SLIDE_W, Inches(0.04), fill_color=C_ACCENT)


def make_content_slide(prs, title, build_fn):
    """创建内容页：标题栏 + 内容构建函数"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    make_title_bar(slide, title)
    build_fn(slide)
    return slide


def make_chart_slide(prs, chart_title, chart_path):
    """创建图件专页：标题 + 图片居中，无长文"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(slide, chart_title)

    if chart_path.exists():
        # 计算图片尺寸，适配页面（留边距）
        max_w = Inches(12.0)
        max_h = Inches(5.8)
        from PIL import Image
        with Image.open(str(chart_path)) as img:
            img_w, img_h = img.size
        ratio = min(max_w / Emu(int(img_w * 914400 / 96)),
                    max_h / Emu(int(img_h * 914400 / 96)))
        pic_w = int(img_w * 914400 / 96 * ratio)
        pic_h = int(img_h * 914400 / 96 * ratio)
        left = int((SLIDE_W - pic_w) / 2)
        top = Inches(1.3)
        slide.shapes.add_picture(str(chart_path), left, top, pic_w, pic_h)
    else:
        add_textbox(slide, Inches(2), Inches(3), Inches(9), Inches(1),
                    text=f"[图件缺失: {chart_path.name}]",
                    font_size=18, font_color=C_RED, alignment=PP_ALIGN.CENTER)
    return slide


# ── 封面页 ──

def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 深蓝全背景
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_DARK)
    # 中部装饰线
    add_rect(slide, Inches(1.5), Inches(2.6), Inches(10.3), Inches(0.03), fill_color=C_ACCENT)
    # 主标题
    add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.2),
                text="执法音视频集中存储与管理系统", font_size=40,
                font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # 副标题
    add_textbox(slide, Inches(1.5), Inches(3.9), Inches(10.3), Inches(0.8),
                text="建 设 方 案 汇 报", font_size=30,
                font_color=RGBColor(0xAA, 0xBB, 0xDD), bold=False,
                alignment=PP_ALIGN.CENTER)
    # 底部装饰线
    add_rect(slide, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.03), fill_color=C_ACCENT)
    # 底部信息
    add_textbox(slide, Inches(1.5), Inches(5.2), Inches(10.3), Inches(0.5),
                text="文件状态：交付稿    密级等级：内部    版本号：V1.0",
                font_size=14, font_color=RGBColor(0x88, 0x99, 0xBB),
                alignment=PP_ALIGN.CENTER)


# ── 目录页 ──

def build_toc(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    make_title_bar(slide, "汇报提纲")

    toc_items = [
        "一、项目背景与建设必要性",
        "二、建设目标与原则",
        "三、系统定位与边界",
        "四、总体架构",
        "五、数据与存储策略",
        "六、部署结构",
        "七、投资与预算",
        "八、风险与控制",
        "九、建设成效与结论",
    ]
    tf = add_textbox(slide, Inches(2), Inches(1.5), Inches(9), Inches(5.5),
                     text="", font_size=20)
    for item in toc_items:
        add_paragraph(tf, item, font_size=22, font_color=C_TEXT, bold=False,
                      space_before=Pt(10), space_after=Pt(6))


# ── 第1页：项目背景与建设必要性 ──

def build_p1(slide):
    # 左栏：背景
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.5),
                text="项目背景", font_size=20, font_color=C_ACCENT, bold=True)
    tf = add_textbox(slide, Inches(0.5), Inches(1.9), Inches(5.8), Inches(2.2),
                     text="", font_size=13)
    add_bullet(tf, "执法音视频数据量持续增长，现有存储体系已积累约 300TB 数据", font_size=13)
    add_bullet(tf, "日新增约 0.6TB，年新增 200TB 以上，现有系统面临容量瓶颈", font_size=13)
    add_bullet(tf, "执法音视频数据属于执法证据数据，对完整性与不可篡改性具有法定要求", font_size=13)
    add_bullet(tf, "信息技术自主可控要求在存储基础设施层面落实国产化建设", font_size=13)

    # 分隔线
    add_rect(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)

    # 下栏：建设必要性
    add_textbox(slide, Inches(0.5), Inches(4.4), Inches(5.8), Inches(0.5),
                text="建设必要性", font_size=20, font_color=C_ACCENT, bold=True)

    # 四个必要性卡片
    cards = [
        ("存储容量不足", "年新增200TB+\n在线周期≥3年\n现有体系无法承载"),
        ("数据管理能力不足", "缺乏全生命周期管理\n检索效率低\n调阅能力有限"),
        ("证据保护能力不足", "缺乏完整性校验\n无防篡改机制\n证据安全无保障"),
        ("国产化要求", "存储基础设施层\n服务器硬件平台\n操作系统层"),
    ]
    for i, (title, desc) in enumerate(cards):
        x = Inches(0.5 + i * 3.15)
        add_rect(slide, x, Inches(5.0), Inches(2.9), Inches(2.0), fill_color=C_LIGHT)
        add_textbox(slide, x + Inches(0.15), Inches(5.1), Inches(2.6), Inches(0.4),
                    text=title, font_size=14, font_color=C_DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(5.5), Inches(2.6), Inches(1.3),
                    text=desc, font_size=12, font_color=C_GRAY,
                    alignment=PP_ALIGN.CENTER)


# ── 第2页：建设目标与原则 ──

def build_p2(slide):
    # 左栏：建设目标
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.5),
                text="建设目标", font_size=20, font_color=C_ACCENT, bold=True)
    tf = add_textbox(slide, Inches(0.5), Inches(1.9), Inches(6), Inches(2.8),
                     text="", font_size=13)
    add_paragraph(tf, "业务目标", font_size=15, font_color=C_DARK, bold=True, space_before=Pt(2))
    add_bullet(tf, "建立统一的执法音视频数据集中存储体系", font_size=13)
    add_bullet(tf, "有效可用存储容量 ≥700TB，在线存储周期 ≥3 年", font_size=13)
    add_bullet(tf, "建立检索调阅服务能力，支撑执法监督与规范化管理", font_size=13)
    add_bullet(tf, "建立数据全生命周期管理机制", font_size=13)
    add_paragraph(tf, "技术目标", font_size=15, font_color=C_DARK, bold=True, space_before=Pt(8))
    add_bullet(tf, "Scale-Out NAS 横向扩展架构，统一文件存储资源池", font_size=13)
    add_bullet(tf, "国产化：服务器 + 操作系统 + 达梦数据库 V8.4", font_size=13)
    add_bullet(tf, "公安专网部署，削峰与分时上传机制", font_size=13)

    # 右栏：建设原则
    add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.5),
                text="建设原则", font_size=20, font_color=C_ACCENT, bold=True)
    principles = [
        ("集约化", "统一存储资源池，各接入通道\n数据汇聚至同一资源池"),
        ("稳定可靠", "双机部署+RAID保护+冗余链路\n单点故障不中断服务"),
        ("横向扩展", "增加NAS节点在线扩容\n对应用层透明，不中断服务"),
        ("自主可控", "国产CPU/OS + 达梦V8.4\n不涉及既有系统改造"),
    ]
    for i, (title, desc) in enumerate(principles):
        y = Inches(2.0 + i * 1.3)
        add_rect(slide, Inches(7), y, Inches(5.5), Inches(1.1), fill_color=C_LIGHT)
        add_textbox(slide, Inches(7.2), y + Inches(0.05), Inches(1.8), Inches(1.0),
                    text=title, font_size=15, font_color=C_DARK, bold=True)
        add_textbox(slide, Inches(9.2), y + Inches(0.05), Inches(3.1), Inches(1.0),
                    text=desc, font_size=12, font_color=C_GRAY)


# ── 第3页：系统定位与边界 ──

def build_p3(slide):
    # 系统定位
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.5),
                text="系统定位", font_size=20, font_color=C_ACCENT, bold=True)
    add_textbox(slide, Inches(0.5), Inches(1.9), Inches(12), Inches(0.8),
                text="执法音视频数据存储与管理支撑系统\n核心职责：数据承载、集中存储、检索调阅、监督支撑",
                font_size=16, font_color=C_TEXT)

    add_rect(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)

    # 建设范围 vs 不包含范围
    add_textbox(slide, Inches(0.5), Inches(3.1), Inches(5.8), Inches(0.4),
                text="建设范围", font_size=18, font_color=RGBColor(0x19, 0x87, 0x54), bold=True)
    tf_in = add_textbox(slide, Inches(0.5), Inches(3.6), Inches(5.8), Inches(3.2),
                        text="", font_size=13)
    for item in [
        "执法音视频接入节点",
        "Scale-Out NAS 集中存储资源池",
        "管理与检索平台（达梦数据库 V8.4）",
        "访问与调阅终端",
        "网络交换与传输链路",
    ]:
        add_bullet(tf_in, item, font_size=13, font_color=C_TEXT)

    add_textbox(slide, Inches(7), Inches(3.1), Inches(5.8), Inches(0.4),
                text="不包含范围", font_size=18, font_color=C_RED, bold=True)
    tf_out = add_textbox(slide, Inches(7), Inches(3.6), Inches(5.8), Inches(3.2),
                         text="", font_size=13)
    for item in [
        "不涉及既有执法业务系统功能改造",
        "不涉及执法终端设备更新或替换",
        "不涉及执法流程制度调整",
        "不包含历史系统替换或全量数据迁移",
        "不承担实时视频监控/调度/流媒体分发",
        "不涉及双中心/异地容灾/云化部署",
    ]:
        add_bullet(tf_out, item, font_size=13, font_color=C_GRAY)

    # 数据接入范围
    add_rect(slide, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)
    add_textbox(slide, Inches(0.5), Inches(6.5), Inches(2.5), Inches(0.4),
                text="数据接入范围：", font_size=14, font_color=C_DARK, bold=True)
    add_textbox(slide, Inches(3.0), Inches(6.5), Inches(9.5), Inches(0.4),
                text="5G执法记录仪回传 | 采集站导入 | 执法仪本地导入 | 既有平台归集",
                font_size=14, font_color=C_TEXT)


# ── 第4页 placeholder（总体架构文字说明） ──

def build_p4(slide):
    # 左栏：系统组成
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4),
                text="系统组成", font_size=18, font_color=C_ACCENT, bold=True)
    tf = add_textbox(slide, Inches(0.5), Inches(1.8), Inches(6), Inches(3.2),
                     text="", font_size=13)
    components = [
        ("接入节点", "多源数据唯一入口，完整性校验后写入存储"),
        ("管理与检索平台", "应用服务器部署，NFS/SMB挂载存储池，JDBC连接达梦"),
        ("NAS 存储资源池", "统一文件存储，唯一数据物理承载层"),
        ("访问终端", "通过平台调阅，不直接访问存储资源池"),
        ("网络交换", "业务网络与存储网络(≥10GbE)逻辑隔离"),
    ]
    for title, desc in components:
        add_bullet(tf, f"{title}：{desc}", font_size=12)

    # 右栏：数据路径
    add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.4),
                text="数据路径", font_size=18, font_color=C_ACCENT, bold=True)
    tf2 = add_textbox(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(3.2),
                      text="", font_size=13)
    add_paragraph(tf2, "写入路径", font_size=14, font_color=C_DARK, bold=True, space_before=Pt(2))
    add_bullet(tf2, "接入节点 → 校验 → 存储网络写入NAS → 数据库建立元数据", font_size=12)
    add_paragraph(tf2, "调阅路径", font_size=14, font_color=C_DARK, bold=True, space_before=Pt(8))
    add_bullet(tf2, "终端 → 平台查询数据库 → 存储网络读取NAS → 返回播放流", font_size=12)
    add_paragraph(tf2, "运维路径", font_size=14, font_color=C_DARK, bold=True, space_before=Pt(8))
    add_bullet(tf2, "存储管理服务器采集状态 → 平台汇聚监控 → 运维终端展示", font_size=12)

    add_rect(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)

    # 底部：部署边界
    add_textbox(slide, Inches(0.5), Inches(5.4), Inches(12), Inches(0.4),
                text="部署边界", font_size=16, font_color=C_DARK, bold=True)
    tf3 = add_textbox(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(1.0),
                      text="", font_size=12)
    add_bullet(tf3, "中心化部署，全部设备集中于公安专网机房", font_size=12)
    add_bullet(tf3, "不涉及双中心/异地容灾/云化部署/跨区域调度，不依赖外部云平台或第三方服务", font_size=12)

    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                text="（详见下页：系统总体架构图）", font_size=14,
                font_color=C_ACCENT, alignment=PP_ALIGN.CENTER)


# ── 第5页 placeholder（数据与存储策略文字说明） ──

def build_p5(slide):
    # 左栏：存储架构
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(0.4),
                text="存储架构", font_size=18, font_color=C_ACCENT, bold=True)
    tf = add_textbox(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.5),
                     text="", font_size=13)
    add_bullet(tf, "Scale-Out NAS 横向扩展架构", font_size=13, bold=True)
    add_bullet(tf, "统一文件存储资源池，NFS/SMB 标准接口", font_size=13)
    add_bullet(tf, "单节点有效容量 ≥810TB（45×18TB）", font_size=13)
    add_bullet(tf, "节点冗余 + RAID 等价数据保护", font_size=13)
    add_bullet(tf, "节点互联带宽 ≥10GbE", font_size=13)

    # 右栏：数据策略
    add_textbox(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(0.4),
                text="数据保护策略", font_size=18, font_color=C_ACCENT, bold=True)
    tf2 = add_textbox(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(2.5),
                      text="", font_size=13)
    add_bullet(tf2, "原始证据数据完整保存，不可修改", font_size=13, bold=True)
    add_bullet(tf2, "全流程完整性校验机制", font_size=13)
    add_bullet(tf2, "压缩/转码仅作用于归档副本（≤20%）", font_size=13)
    add_bullet(tf2, "压缩数据不作为证据载体", font_size=13)

    add_rect(slide, Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)

    # 关键指标
    add_textbox(slide, Inches(0.5), Inches(4.7), Inches(12), Inches(0.4),
                text="关键数据指标", font_size=18, font_color=C_ACCENT, bold=True)
    metrics = [
        ("日新增", "≈0.6TB"),
        ("年新增", "≥200TB"),
        ("有效容量", "≥700TB"),
        ("在线周期", "≥3 年"),
        ("并发用户", "≥20"),
        ("并发播放", "≥20 路"),
    ]
    for i, (label, value) in enumerate(metrics):
        x = Inches(0.5 + i * 2.1)
        add_rect(slide, x, Inches(5.2), Inches(1.9), Inches(1.5), fill_color=C_LIGHT)
        add_textbox(slide, x, Inches(5.3), Inches(1.9), Inches(0.5),
                    text=value, font_size=22, font_color=C_DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, Inches(5.9), Inches(1.9), Inches(0.5),
                    text=label, font_size=13, font_color=C_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                text="（详见下页：数据流转图）", font_size=14,
                font_color=C_ACCENT, alignment=PP_ALIGN.CENTER)


# ── 第6页 placeholder（部署结构文字说明） ──

def build_p6(slide):
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
                text="中心化部署 — 公安专网机房", font_size=18,
                font_color=C_ACCENT, bold=True)

    # 设备清单卡片
    devices = [
        ("应用服务器 ×2", "管理与检索平台\n数据接入服务\n国产CPU/OS"),
        ("存储服务器 ×1", "Scale-Out NAS\n有效容量 ≥810TB\nNFS/SMB 接口"),
        ("存储管理服务器 ×1", "存储管理服务\n运维监控\n国产CPU/OS"),
        ("达梦数据库 ×1", "元数据/业务数据\n索引/审计日志\n国产数据库 V8.4"),
    ]
    for i, (title, desc) in enumerate(devices):
        x = Inches(0.5 + i * 3.15)
        add_rect(slide, x, Inches(2.0), Inches(2.9), Inches(2.5), fill_color=C_LIGHT)
        add_textbox(slide, x + Inches(0.15), Inches(2.1), Inches(2.6), Inches(0.5),
                    text=title, font_size=15, font_color=C_DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), Inches(2.6), Inches(2.6), Inches(1.7),
                    text=desc, font_size=13, font_color=C_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)

    # 部署约束
    add_textbox(slide, Inches(0.5), Inches(5.0), Inches(12), Inches(0.4),
                text="部署约束", font_size=18, font_color=C_ACCENT, bold=True)
    tf = add_textbox(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
                     text="", font_size=13)
    add_bullet(tf, "系统采用中心化部署模式，全部设备集中部署于公安专网机房", font_size=13)
    add_bullet(tf, "不涉及双中心架构、异地容灾、云化部署及跨区域数据调度", font_size=13)
    add_bullet(tf, "新建系统独立部署，不影响既有系统运行", font_size=13)
    add_bullet(tf, "既有系统（约300TB历史数据）保留运行，承担历史查询与调阅", font_size=13)


# ── 第7页：投资与预算 ──

def build_p7(slide):
    # 总投资金额
    add_rect(slide, Inches(4), Inches(1.3), Inches(5.3), Inches(1.0), fill_color=C_LIGHT)
    add_textbox(slide, Inches(4), Inches(1.35), Inches(5.3), Inches(0.5),
                text="项目投资总额", font_size=16, font_color=C_GRAY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(4), Inches(1.7), Inches(5.3), Inches(0.6),
                text="107.4 万元", font_size=30, font_color=C_DARK, bold=True,
                alignment=PP_ALIGN.CENTER)

    # 三个分类
    cats = [
        ("软硬件购置费", "81 万元",
         "应用服务器 ×2    22万\n存储服务器 ×1    45万\n存储管理服务器 ×1  5万\n达梦数据库 ×1     9万"),
        ("应用软件费", "20 万元",
         "执法音视频管理与\n存储应用软件 ×1"),
        ("其它费用", "6.4 万元",
         "系统集成费  6.4万\n监理费      0万\n等保测评    0万"),
    ]
    for i, (title, amount, detail) in enumerate(cats):
        x = Inches(0.5 + i * 4.2)
        add_rect(slide, x, Inches(2.8), Inches(3.8), Inches(4.0), fill_color=C_LIGHT)
        add_textbox(slide, x + Inches(0.2), Inches(2.9), Inches(3.4), Inches(0.5),
                    text=title, font_size=18, font_color=C_DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(3.4), Inches(3.4), Inches(0.5),
                    text=amount, font_size=22, font_color=C_ACCENT, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), Inches(4.0), Inches(3.4), Inches(2.5),
                    text=detail, font_size=13, font_color=C_GRAY,
                    alignment=PP_ALIGN.LEFT)

    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
                text="资金来源：财政拨款    |    各分项之和 = 总计 107.4 万元",
                font_size=12, font_color=C_GRAY, alignment=PP_ALIGN.CENTER)


# ── 第8页：风险与控制 ──

def build_p8(slide):
    risks = [
        ("设备供货风险", "设备交付延迟或到货不合格",
         "合同明确交付时限与违约责任\n到货后立即预验收与兼容测试"),
        ("既有系统影响", "新建系统部署影响既有系统运行",
         "独立部署，不共享存储/应用/数据库\n部署过程不对既有系统变更"),
        ("数据安全风险", "试运行期间数据丢失或损坏",
         "既有系统保持运行保障连续性\n新增数据在新系统完成校验入库"),
        ("网络冲击风险", "大批量接入对专网带宽冲击",
         "削峰与分时上传机制\n试运行初期控制接入量逐步放量"),
        ("集成对接风险", "与既有执法平台接口对接失败",
         "优先接口对接验证\n标准协议 + 预留联调时间"),
        ("证据完整性", "传输/存储中数据损坏篡改",
         "全流程哈希校验机制\nRAID等价保护 + 审计追溯"),
    ]
    for i, (title, risk, measure) in enumerate(risks):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.3 + row * 3.0)
        add_rect(slide, x, y, Inches(3.8), Inches(2.6), fill_color=C_LIGHT)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.4),
                    text=title, font_size=16, font_color=C_DARK, bold=True)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(0.8),
                    text=f"风险：{risk}", font_size=12, font_color=C_RED)
        add_textbox(slide, x + Inches(0.15), y + Inches(1.3), Inches(3.5), Inches(1.2),
                    text=f"应对：{measure}", font_size=12, font_color=RGBColor(0x19, 0x87, 0x54))


# ── 第9页：建设成效与结论 ──

def build_p9(slide):
    add_textbox(slide, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
                text="预期建设成效", font_size=20, font_color=C_ACCENT, bold=True)

    achievements = [
        ("存储能力", "≥700TB 有效容量\n满足 ≥3 年在线存储"),
        ("接入能力", "多源统一接入\n日均 ≥0.6TB 写入"),
        ("服务能力", "≥20 并发用户\n≥20 路并发播放"),
        ("管理能力", "全生命周期自动化\n人工干预显著减少"),
        ("安全能力", "证据完整性校验 100%\n全操作审计追溯"),
        ("自主可控", "核心基础设施国产化\n达梦数据库 V8.4"),
    ]
    for i, (title, desc) in enumerate(achievements):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.9 + row * 2.0)
        add_rect(slide, x, y, Inches(3.8), Inches(1.7), fill_color=C_LIGHT)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(3.5), Inches(0.4),
                    text=title, font_size=16, font_color=C_DARK, bold=True,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.15), y + Inches(0.55), Inches(3.5), Inches(1.0),
                    text=desc, font_size=14, font_color=C_GRAY,
                    alignment=PP_ALIGN.CENTER)

    add_rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.02), fill_color=C_ACCENT)

    add_textbox(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(1.0),
                text="本项目聚焦执法音视频数据存储与管理基础设施建设，建设范围明确、目标清晰、技术可行、风险可控。\n系统定位为数据存储支撑层，不扩展业务功能、不替代既有系统、不改变执法流程。",
                font_size=15, font_color=C_TEXT, alignment=PP_ALIGN.CENTER)


# ── 结尾页 ──

def build_end(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=C_DARK)
    add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10.3), Inches(1.0),
                text="汇 报 完 毕", font_size=44,
                font_color=C_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.8),
                text="敬 请 指 导", font_size=28,
                font_color=RGBColor(0xAA, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)


# ── 主构建流程 ──

def build():
    prs = new_prs()

    # 封面
    build_cover(prs)
    # 目录
    build_toc(prs)

    # 1. 项目背景与建设必要性
    make_content_slide(prs, "一、项目背景与建设必要性", build_p1)
    # 2. 建设目标与原则
    make_content_slide(prs, "二、建设目标与原则", build_p2)
    # 3. 系统定位与边界
    make_content_slide(prs, "三、系统定位与边界", build_p3)
    # 4. 总体架构（文字）
    make_content_slide(prs, "四、总体架构", build_p4)
    # 4-图：系统总体架构图
    make_chart_slide(prs, "系统总体架构图", CHARTS["系统总体架构图"])

    # 5. 数据与存储策略（文字）
    make_content_slide(prs, "五、数据与存储策略", build_p5)
    # 5-图：数据流转图
    make_chart_slide(prs, "数据流转图", CHARTS["数据流转图"])

    # 6. 部署结构（文字）
    make_content_slide(prs, "六、部署结构", build_p6)
    # 6-图：网络拓扑图（独占一页）
    make_chart_slide(prs, "网络拓扑图", CHARTS["网络拓扑图"])
    # 6-图：部署结构图（独占一页）
    make_chart_slide(prs, "部署结构图", CHARTS["部署结构图"])

    # 7. 投资与预算
    make_content_slide(prs, "七、投资与预算", build_p7)
    # 8. 风险与控制
    make_content_slide(prs, "八、风险与控制", build_p8)
    # 9. 建设成效与结论
    make_content_slide(prs, "九、建设成效与结论", build_p9)

    # 结尾页
    build_end(prs)

    prs.save(str(OUTPUT))
    print(f"Done: {OUTPUT}")


if __name__ == '__main__':
    build()
