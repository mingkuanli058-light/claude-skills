# -*- coding: utf-8 -*-
"""
Proposal Writing Skill 培训 PPT 生成脚本
政务风格：深蓝色顶栏 + 白色主体，专业克制
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 颜色常量 ──────────────────────────────────────────────
C_NAVY      = RGBColor(0x1B, 0x2D, 0x4F)   # 深海蓝（主色）
C_BLUE      = RGBColor(0x1F, 0x4E, 0x79)   # 钢蓝（副色）
C_ACCENT    = RGBColor(0x2E, 0x74, 0xB5)   # 强调蓝
C_LIGHT_BG  = RGBColor(0xF2, 0xF5, 0xF9)   # 极浅蓝背景
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY_TEXT = RGBColor(0x44, 0x44, 0x55)
C_GRAY_LINE = RGBColor(0xCC, 0xD6, 0xE0)
C_WARN      = RGBColor(0xC0, 0x39, 0x2B)   # 红色（禁止项）
C_OK        = RGBColor(0x1A, 0x6B, 0x3C)   # 绿色（允许项）
C_TABLE_H   = RGBColor(0x1B, 0x2D, 0x4F)   # 表头
C_TABLE_R1  = RGBColor(0xF2, 0xF5, 0xF9)   # 表格偶数行
C_TABLE_R2  = RGBColor(0xFF, 0xFF, 0xFF)   # 表格奇数行
C_YELLOW    = RGBColor(0xF5, 0xA6, 0x23)   # 金色标注


# ── 幻灯片尺寸 16:9 ───────────────────────────────────────
W = Cm(33.867)
H = Cm(19.05)

FONT_CN = "Microsoft YaHei"   # 正文中文
FONT_TITLE = "Microsoft YaHei"


def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)


# ── 基础绘图工具 ──────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text,
                font_size=18, bold=False, color=C_DARK_TEXT,
                align=PP_ALIGN.LEFT, font=FONT_CN,
                line_spacing=None, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    if line_spacing:
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
        spcPts.set('val', str(int(line_spacing * 100)))
    return txb


def set_run(p, text, font_size=16, bold=False, color=C_DARK_TEXT,
            font=FONT_CN, italic=False):
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return run


def add_multiline_textbox(slide, x, y, w, h, lines, font_size=15,
                           color=C_DARK_TEXT, font=FONT_CN,
                           bold_first=False, spacing_before=6):
    "\"\"lines: list of (text, bold, color_override)\"\""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            text, bold, col = item, False, color
        else:
            text = item[0]
            bold = item[1] if len(item) > 1 else False
            col  = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing_before if i > 0 else 0)
        run = p.add_run()
        run.text = text
        run.font.name  = font
        run.font.size  = Pt(font_size)
        run.font.bold  = bold and (not bold_first or i == 0)
        if bold_first and i == 0:
            run.font.bold = True
        run.font.color.rgb = col
    return txb


# ── 顶栏 + 分隔线（标准页眉） ─────────────────────────────

def draw_header(slide, title, subtitle=None):
    "\"\"顶栏：深蓝背景 + 白色标题\"\""
    bar_h = Cm(2.5)
    add_rect(slide, 0, 0, W, bar_h, C_NAVY)

    # 左侧竖线装饰
    add_rect(slide, Cm(0.8), Cm(0.4), Cm(0.18), Cm(1.7), C_YELLOW)

    # 标题
    add_textbox(slide, Cm(1.3), Cm(0.35), Cm(28), Cm(1.8),
                title, font_size=22, bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)

    # 副标题
    if subtitle:
        add_textbox(slide, Cm(1.3), Cm(1.8), Cm(28), Cm(0.6),
                    subtitle, font_size=11, bold=False,
                    color=RGBColor(0xAA, 0xBB, 0xCC),
                    align=PP_ALIGN.LEFT)

    # 底部横线
    add_rect(slide, 0, bar_h, W, Cm(0.06), C_ACCENT)


def draw_page_num(slide, num, total=20):
    txt = f"{num:02d} / {total:02d}"
    add_textbox(slide, W - Cm(3.5), H - Cm(0.7), Cm(3.2), Cm(0.6),
                txt, font_size=10, color=C_GRAY_TEXT,
                align=PP_ALIGN.RIGHT)


def draw_footer_line(slide):
    add_rect(slide, 0, H - Cm(0.75), W, Cm(0.05), C_GRAY_LINE)


# ── 表格绘制工具 ──────────────────────────────────────────

def draw_table(slide, x, y, w, col_widths, rows_data,
               font_size=13, row_h=Cm(0.72)):
    """
    rows_data: list of lists of str
    col_widths: list of Cm values, must sum to w
    row 0 = header (navy bg, white text)
    """
    from pptx.util import Emu
    n_rows = len(rows_data)
    n_cols = len(col_widths)

    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w,
                                  row_h * n_rows).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw

    for ri, row in enumerate(rows_data):
        tbl.rows[ri].height = row_h
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            tf = cell.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.space_before = Pt(2)
                for run in para.runs:
                    run.font.name  = FONT_CN
                    run.font.size  = Pt(font_size)
                    if ri == 0:
                        run.font.bold  = True
                        run.font.color.rgb = C_WHITE
                    else:
                        run.font.color.rgb = C_DARK_TEXT
            # 背景色
            fill = cell.fill
            fill.solid()
            if ri == 0:
                fill.fore_color.rgb = C_TABLE_H
            elif ri % 2 == 0:
                fill.fore_color.rgb = C_TABLE_R1
            else:
                fill.fore_color.rgb = C_TABLE_R2

    return tbl


# ── 带标签的内容块 ────────────────────────────────────────

def draw_label_block(slide, x, y, w, label, items,
                     label_color=C_BLUE, font_size=14,
                     item_prefix="▸  "):
    "\"\"左侧竖线色块 + 标签 + 要点列表\"\""
    add_rect(slide, x, y, Cm(0.15), Cm(0.5 + len(items) * 0.65),
             label_color)
    add_textbox(slide, x + Cm(0.35), y, w, Cm(0.5),
                label, font_size=13, bold=True, color=label_color)
    for i, item in enumerate(items):
        add_textbox(slide, x + Cm(0.35), y + Cm(0.5 + i * 0.65),
                    w, Cm(0.6),
                    item_prefix + item, font_size=font_size,
                    color=C_DARK_TEXT)


# ══════════════════════════════════════════════════════════
# 各页生成函数
# ══════════════════════════════════════════════════════════

def slide_01_cover(prs):
    "\"\"封面\"\""
    slide = blank_slide(prs)

    # 全深蓝背景上半
    add_rect(slide, 0, 0, W, H * 0.55, C_NAVY)
    # 下半浅色
    add_rect(slide, 0, H * 0.55, W, H * 0.45, C_LIGHT_BG)

    # 装饰竖线
    add_rect(slide, Cm(1.5), Cm(3.2), Cm(0.25), Cm(4.0), C_YELLOW)

    # 主标题
    add_textbox(slide, Cm(2.2), Cm(3.0), Cm(26), Cm(2.2),
                "Proposal Writing Skill",
                font_size=38, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(2.2), Cm(5.0), Cm(26), Cm(1.2),
                "方案生产过程控制体系　内部培训材料",
                font_size=20, bold=False,
                color=RGBColor(0xAA, 0xBB, 0xCC))

    # 下方信息
    add_rect(slide, Cm(2.2), Cm(11.5), Cm(20), Cm(0.06), C_ACCENT)
    add_textbox(slide, Cm(2.2), Cm(11.8), Cm(25), Cm(0.7),
                "适用岗位：技术负责人 · 售前工程师 · 投标专员",
                font_size=14, color=C_GRAY_TEXT)
    add_textbox(slide, Cm(2.2), Cm(12.6), Cm(25), Cm(0.7),
                "预计时长：30 – 40 分钟    共 20 页",
                font_size=14, color=C_GRAY_TEXT)
    add_textbox(slide, Cm(2.2), Cm(13.4), Cm(25), Cm(0.7),
                "2026 · 内部使用，请勿对外发布",
                font_size=13, color=C_GRAY_TEXT)

    draw_page_num(slide, 1)


def slide_02_overview(prs):
    slide = blank_slide(prs)
    draw_header(slide, "培训说明", "Training Overview")
    draw_footer_line(slide)
    draw_page_num(slide, 2)

    y0 = Cm(3.0)
    # 培训目标
    add_textbox(slide, Cm(1.5), y0, Cm(30), Cm(0.6),
                "▌ 本次培训目标", font_size=15, bold=True, color=C_NAVY)
    goals = [
        "理解 Proposal Writing Skill 的定位与职责边界",
        "掌握方案生产十步法的执行顺序与各阶段控制节点",
        "了解技术路线冻结机制与投资控制逻辑",
        "掌握图件体系生成纪律与格式要求",
        "识别并规避常见写作违规行为",
        "使用提交前检查清单进行自查",
    ]
    for i, g in enumerate(goals):
        add_textbox(slide, Cm(2.0), y0 + Cm(0.7 + i * 0.78), Cm(29), Cm(0.65),
                    f"  {i+1}.  {g}", font_size=14, color=C_DARK_TEXT)

    # 适用场景
    y1 = y0 + Cm(0.7 + len(goals) * 0.78 + 0.4)
    add_textbox(slide, Cm(1.5), y1, Cm(30), Cm(0.6),
                "▌ 适用场景", font_size=15, bold=True, color=C_NAVY)
    scenes = [
        "使用 Claude Code + Proposal Writing Skill 生成方案",
        "人工编制方案时的过程控制参考",
        "投标材料合规性自查",
    ]
    for i, s in enumerate(scenes):
        add_textbox(slide, Cm(2.0), y1 + Cm(0.7 + i * 0.72), Cm(29), Cm(0.62),
                    f"  ▸  {s}", font_size=14, color=C_DARK_TEXT)

    # 说明框
    add_rect(slide, Cm(1.5), H - Cm(2.4), Cm(30.9), Cm(1.5),
             C_LIGHT_BG)
    add_textbox(slide, Cm(2.0), H - Cm(2.3), Cm(30), Cm(1.3),
                "本培训不讨论技术路线选型，不评价设计方案优劣，\n仅针对\"方案如何生产\"这一工程控制问题。",
                font_size=13, color=C_GRAY_TEXT)


def slide_03_positioning(prs):
    slide = blank_slide(prs)
    draw_header(slide, "方法定位", "What This Skill Is / Is Not")
    draw_footer_line(slide)
    draw_page_num(slide, 3)

    # 左栏：不是什么
    add_rect(slide, Cm(1.2), Cm(3.0), Cm(14.5), Cm(7.5),
             RGBColor(0xFD, 0xF2, 0xF2))
    add_rect(slide, Cm(1.2), Cm(3.0), Cm(14.5), Cm(1.0), C_WARN)
    add_textbox(slide, Cm(1.5), Cm(3.05), Cm(14), Cm(0.85),
                "✕  不是什么", font_size=14, bold=True, color=C_WHITE)
    not_items = [
        '"AI 写方案"工具',
        "技术顾问",
        "提供技术建议的角色",
        "参与需求判断的角色",
    ]
    for i, t in enumerate(not_items):
        add_textbox(slide, Cm(1.8), Cm(4.2 + i * 0.85), Cm(13.5), Cm(0.75),
                    f"  ✕  {t}", font_size=14, color=C_WARN)

    # 右栏：是什么
    add_rect(slide, Cm(17.2), Cm(3.0), Cm(15.0), Cm(7.5),
             RGBColor(0xF0, 0xF7, 0xFF))
    add_rect(slide, Cm(17.2), Cm(3.0), Cm(15.0), Cm(1.0), C_BLUE)
    add_textbox(slide, Cm(17.5), Cm(3.05), Cm(14.5), Cm(0.85),
                "✔  是什么", font_size=14, bold=True, color=C_WHITE)
    is_items = [
        "政府 / 公安信息化项目方案编制过程控制体系",
        "控制模型写作行为与输出结构的执行工具",
        "强制模型服从权威文件、禁止自由发挥",
    ]
    for i, t in enumerate(is_items):
        add_textbox(slide, Cm(17.7), Cm(4.2 + i * 0.9), Cm(14.2), Cm(0.8),
                    f"  ✔  {t}", font_size=14, color=C_BLUE)

    # 核心定义框
    add_rect(slide, Cm(1.2), Cm(11.1), Cm(31.0), Cm(2.0), C_NAVY)
    add_textbox(slide, Cm(1.8), Cm(11.25), Cm(30), Cm(0.6),
                "核心定义", font_size=12, bold=True,
                color=RGBColor(0xAA, 0xBB, 0xCC))
    add_textbox(slide, Cm(1.8), Cm(11.8), Cm(30), Cm(1.0),
                "方案的本质 = 决策与约束的表达载体，而非文本创作行为。",
                font_size=16, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(1.8), H - Cm(1.6), Cm(30), Cm(0.6),
                "Skill 负责\"怎么说\"，不负责\"说什么\"。\"说什么\"由 specs/ 权威文件决定。",
                font_size=12, color=C_GRAY_TEXT)


def slide_04_principles(prs):
    slide = blank_slide(prs)
    draw_header(slide, "核心原则", "Four Non-Negotiable Priorities")
    draw_footer_line(slide)
    draw_page_num(slide, 4)

    principles = [
        ("1", "需求优先于方案",
         "没有固化需求，不得生成方案。口头共识不能替代结构化文件。"),
        ("2", "决策优先于写作",
         "未完成技术路线裁决，不得进入章节写作阶段。"),
        ("3", "风险控制优先于技术偏好",
         "不以技术先进性作为方案选型的主要依据。"),
        ("4", "投资合理性优先于功能扩展",
         "方案能力边界由投资模型决定，不得超出 investment_schema 范围。"),
    ]

    for i, (num, title, desc) in enumerate(principles):
        x = Cm(1.2) + (i % 2) * Cm(16.3)
        y = Cm(3.2) + (i // 2) * Cm(4.2)
        w = Cm(15.5)

        add_rect(slide, x, y, w, Cm(3.8), C_LIGHT_BG)
        add_rect(slide, x, y, Cm(1.0), Cm(3.8), C_NAVY)
        add_textbox(slide, x + Cm(0.1), y + Cm(0.9), Cm(0.8), Cm(1.0),
                    num, font_size=22, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Cm(1.2), y + Cm(0.4), Cm(13.8), Cm(0.75),
                    title, font_size=16, bold=True, color=C_NAVY)
        add_textbox(slide, x + Cm(1.2), y + Cm(1.3), Cm(13.5), Cm(2.0),
                    desc, font_size=13, color=C_GRAY_TEXT)

    # 底部说明
    add_textbox(slide, Cm(1.2), H - Cm(1.6), Cm(30), Cm(0.65),
                "以上四条原则是评审阶段最常见的被质疑点，也是最需要在编制阶段提前控制的风险维度。",
                font_size=12, color=C_GRAY_TEXT)


def slide_05_file_structure(prs):
    slide = blank_slide(prs)
    draw_header(slide, "整体文件体系架构", "File System Architecture")
    draw_footer_line(slide)
    draw_page_num(slide, 5)

    # 目录树（代码块风格）
    add_rect(slide, Cm(1.2), Cm(2.9), Cm(19.5), Cm(11.8),
             RGBColor(0x1E, 0x1E, 0x2E))
    tree_lines = [
        ("proposal-writing-skill/", True, C_YELLOW),
        ("├── SKILL.md", False, RGBColor(0x7E, 0xC8, 0x50)),
        ("│   ← 唯一规则权威文件", False, C_GRAY_LINE),
        ("├── specs/", True, RGBColor(0x56, 0xB6, 0xC2)),
        ("│   ├── 需求.md          ← 业务事实 SSoT", False, C_WHITE),
        ("│   ├── technical_baseline.md ← 技术路线裁决", False, C_WHITE),
        ("│   ├── architecture_scope.md ← 架构范围约束", False, C_WHITE),
        ("│   ├── investment_schema.md  ← 投资模型权威", False, C_WHITE),
        ("│   ├── chart_registry.md    ← 图件注册清单", False, C_WHITE),
        ("│   └── document_conventions.md", False, C_WHITE),
        ("├── drafts/              ← 草稿区（禁止整体重写）", False,
         RGBColor(0xE0, 0xBF, 0x7A)),
        ("│   └── chart/           ← 推理图（Mermaid）", False,
         RGBColor(0xE0, 0xBF, 0x7A)),
        ("└── output/              ← 正式交付区", False,
         RGBColor(0x98, 0xC3, 0x79)),
        ("    ├── chart/           ← 正式图件（HTML→PNG）", False,
         RGBColor(0x98, 0xC3, 0x79)),
        ("    ├── screenshot/      ← 现网截图（独立管理）", False,
         RGBColor(0x98, 0xC3, 0x79)),
        ("    └── *.md / *.docx    ← 最终方案文本", False,
         RGBColor(0x98, 0xC3, 0x79)),
    ]
    for i, (line, bold, col) in enumerate(tree_lines):
        add_textbox(slide, Cm(1.6), Cm(3.1 + i * 0.66), Cm(19), Cm(0.62),
                    line, font_size=11, bold=bold, color=col,
                    font="Courier New")

    # 右侧规则说明
    rx = Cm(21.3)
    add_textbox(slide, rx, Cm(2.9), Cm(12), Cm(0.65),
                "▌ 关键规则", font_size=14, bold=True, color=C_NAVY)
    rules = [
        ("specs/ 为最高权威层",
         "内容一旦确认不得随意变更，所有决策均以此为依据"),
        ("SKILL.md 是唯一行为规则文件",
         "其余文件均为规则支撑组件，不独立裁决行为"),
        ("drafts/ 内容属于中间态",
         "禁止整体重写，禁止直接作为正式交付物"),
        ("正式图件不得从 drafts/ 复制",
         "必须经过 HTML 渲染后以正式流程生成"),
    ]
    for i, (t, d) in enumerate(rules):
        y = Cm(3.8 + i * 2.2)
        add_rect(slide, rx, y, Cm(12.3), Cm(1.9), C_LIGHT_BG)
        add_rect(slide, rx, y, Cm(0.15), Cm(1.9), C_ACCENT)
        add_textbox(slide, rx + Cm(0.4), y + Cm(0.1), Cm(11.8), Cm(0.6),
                    t, font_size=13, bold=True, color=C_NAVY)
        add_textbox(slide, rx + Cm(0.4), y + Cm(0.7), Cm(11.8), Cm(1.0),
                    d, font_size=12, color=C_GRAY_TEXT)


def slide_06_authority(prs):
    slide = blank_slide(prs)
    draw_header(slide, "权威优先级", "Authority Hierarchy & Conflict Resolution")
    draw_footer_line(slide)
    draw_page_num(slide, 6)

    # 金字塔式优先级
    tiers = [
        (1, "specs/ 目录", "业务事实与约束，所有决策的根基", C_WARN, Cm(10)),
        (2, "drafts/ 目录", "中间草稿，禁止整体重写", C_BLUE, Cm(13)),
        (3, "output/ 目录", "正式交付物", C_ACCENT, Cm(16)),
        (4, "PROPOSAL_PLAN.md", "可丢弃的状态跟踪文件", C_GRAY_TEXT, Cm(19)),
        (5, "模型推理 / 经验判断", "禁止作为事实来源", C_WARN, Cm(22)),
    ]
    for rank, name, desc, col, w in tiers:
        x = (W - w) / 2
        y = Cm(2.7 + (rank - 1) * 1.75)
        h = Cm(1.55)
        add_rect(slide, x, y, w, h, col)
        add_textbox(slide, x + Cm(0.5), y + Cm(0.05), w - Cm(1), Cm(0.65),
                    f"  优先级 {rank}  ·  {name}",
                    font_size=13, bold=True, color=C_WHITE)
        add_textbox(slide, x + Cm(0.5), y + Cm(0.65), w - Cm(1), Cm(0.75),
                    f"  {desc}", font_size=12, color=C_WHITE)

    # 右侧说明
    add_rect(slide, Cm(24.5), Cm(2.7), Cm(8.8), Cm(7.5), C_LIGHT_BG)
    add_textbox(slide, Cm(25.0), Cm(2.9), Cm(8.0), Cm(0.6),
                "specs/ 内部冲突裁决", font_size=13, bold=True, color=C_NAVY)
    add_textbox(slide, Cm(25.0), Cm(3.6), Cm(8.0), Cm(5.5),
                "当 需求.md 与 investment_schema.md\n"
                "在资源规模上存在冲突时：\n\n"
                "▸ investment_schema.md 优先\n\n"
                "需求.md 仅定义业务与功能边界，\n"
                "不裁决资源模型。\n\n"
                "资源配置权属于投资模型文件。",
                font_size=12, color=C_DARK_TEXT)

    add_textbox(slide, Cm(1.5), H - Cm(1.65), Cm(31), Cm(0.65),
                "任何冲突：服从更高权威。不得以正文描述推翻 specs/ 文件。",
                font_size=12, bold=True, color=C_WARN)


def slide_07_tensteps_overview(prs):
    slide = blank_slide(prs)
    draw_header(slide, "方案生产十步法概览", "10-Step Production Workflow")
    draw_footer_line(slide)
    draw_page_num(slide, 7)

    phases = [
        ("需求阶段", [
            ("1", "用户需求输入"),
            ("2", "需求澄清与修正"),
            ("3", "固化为需求.md"),
        ], C_WARN),
        ("决策阶段", [
            ("4", "需求基线评审"),
            ("5A", "技术路线识别"),
            ("5B", "技术路线裁决"),
        ], C_BLUE),
        ("验证阶段", [
            ("6", "原厂询价测算"),
            ("7", "投资模型构建"),
            ("8", "投资合理性评审"),
        ], C_OK),
        ("输出阶段", [
            ("9", "建设方案书写"),
            ("10", "综合评审复核"),
        ], C_NAVY),
    ]

    col_w = Cm(7.8)
    for pi, (phase_name, steps, col) in enumerate(phases):
        px = Cm(1.0) + pi * Cm(8.1)
        # 阶段标题
        add_rect(slide, px, Cm(2.8), col_w, Cm(0.85), col)
        add_textbox(slide, px, Cm(2.85), col_w, Cm(0.75),
                    phase_name, font_size=14, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        # 步骤
        for si, (num, name) in enumerate(steps):
            sy = Cm(3.85 + si * 2.0)
            add_rect(slide, px, sy, col_w, Cm(1.75), C_LIGHT_BG)
            add_rect(slide, px, sy, Cm(1.2), Cm(1.75), col)
            add_textbox(slide, px + Cm(0.05), sy + Cm(0.35), Cm(1.1), Cm(0.85),
                        num, font_size=16, bold=True, color=C_WHITE,
                        align=PP_ALIGN.CENTER)
            add_textbox(slide, px + Cm(1.35), sy + Cm(0.4), Cm(6.1), Cm(0.85),
                        name, font_size=13, bold=True, color=C_DARK_TEXT)
        # 箭头（除最后一组）
        if pi < 3:
            add_textbox(slide, px + col_w + Cm(0.1), Cm(5.0), Cm(0.8), Cm(1.0),
                        "▶", font_size=22, color=C_GRAY_LINE,
                        align=PP_ALIGN.CENTER)

    # 三条禁止规则
    rules = [
        "禁止跳步：任何阶段未完成，不得进入下一阶段",
        "禁止逆序：已固化的文件不得被后续阶段推翻",
        "禁止并行：需求、决策、验证三阶段必须顺序完成",
    ]
    add_rect(slide, Cm(1.0), Cm(13.0), Cm(32.0), Cm(3.5), C_NAVY)
    add_textbox(slide, Cm(1.5), Cm(13.15), Cm(31), Cm(0.55),
                "核心控制逻辑", font_size=12, bold=True,
                color=RGBColor(0xAA, 0xBB, 0xCC))
    for i, r in enumerate(rules):
        add_textbox(slide, Cm(1.8), Cm(13.8 + i * 0.75), Cm(30), Cm(0.65),
                    f"  ✕  {r}", font_size=13, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(1.8), H - Cm(0.9), Cm(30), Cm(0.5),
                "Step 9（方案书写）是整个流程最后进入的环节，不是最先开始的环节。",
                font_size=12, color=C_GRAY_TEXT)


def slide_08_steps123(prs):
    slide = blank_slide(prs)
    draw_header(slide, "需求阶段详解（Step 1–3）", "Requirements Phase")
    draw_footer_line(slide)
    draw_page_num(slide, 8)

    blocks = [
        ("Step 1", "用户需求输入（非结构化阶段）",
         ["来源：会议纪要 / 用户口述 / 问题描述 / 厂商交流材料",
          "信息可能不完整，技术语言可能混乱或自相矛盾",
          "不具备直接写方案的条件",
          "⚠ 禁止行为：Step 1 完成前进入方案写作"],
         C_WARN),
        ("Step 2", "需求澄清与修正（认知治理阶段）",
         ["目标：去除歧义与冲突，剔除不可实施设想，明确系统边界",
          "本阶段属于认知治理，而非技术设计",
          "输出的是\"确认过的事实\"，不是\"设计方案\""],
         C_BLUE),
        ("Step 3", "固化为 specs/需求.md（事实基线阶段）",
         ["必须包含\"现网设备清单\"章节 → 若缺失，禁止进入方案生成",
          "现网结构表达仅可来源于此文件",
          "文件一旦确认，不应随意变更",
          "★ 是后续全部决策与方案生成的唯一事实依据（SSoT）"],
         C_NAVY),
    ]
    for i, (step, title, items, col) in enumerate(blocks):
        y = Cm(3.0 + i * 4.5)
        add_rect(slide, Cm(1.2), y, Cm(32.0), Cm(4.2), C_LIGHT_BG)
        add_rect(slide, Cm(1.2), y, Cm(2.0), Cm(4.2), col)
        add_textbox(slide, Cm(1.25), y + Cm(0.5), Cm(1.9), Cm(0.7),
                    step, font_size=12, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(1.3), y + Cm(1.15), Cm(1.8), Cm(0.8),
                    f"{i+1}", font_size=30, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(3.6), y + Cm(0.25), Cm(29), Cm(0.7),
                    title, font_size=15, bold=True, color=C_NAVY)
        for j, item in enumerate(items):
            c = C_WARN if item.startswith("⚠") or item.startswith("★") else C_DARK_TEXT
            add_textbox(slide, Cm(3.8), y + Cm(1.1 + j * 0.73), Cm(28.5), Cm(0.65),
                        f"  ▸  {item}", font_size=13, color=c)


def slide_09_steps45(prs):
    slide = blank_slide(prs)
    draw_header(slide, "决策阶段详解（Step 4–5B）", "Decision Phase")
    draw_footer_line(slide)
    draw_page_num(slide, 9)

    # Step 4
    add_rect(slide, Cm(1.2), Cm(3.0), Cm(32.0), Cm(3.2), C_LIGHT_BG)
    add_rect(slide, Cm(1.2), Cm(3.0), Cm(1.5), Cm(3.2), C_BLUE)
    add_textbox(slide, Cm(1.25), Cm(3.5), Cm(1.4), Cm(0.7),
                "4", font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3.1), Cm(3.1), Cm(29), Cm(0.7),
                "Step 4 ｜ 需求基线评审（合规校验）", font_size=15, bold=True, color=C_NAVY)
    items4 = ["验证是否存在逻辑冲突", "验证是否符合评审语境（财政、招采合规）",
              "验证是否具备实施与采购可行性",
              "产出：通过评审的 specs/需求.md（确认版本）"]
    for j, t in enumerate(items4):
        add_textbox(slide, Cm(3.3), Cm(3.85 + j * 0.55), Cm(28.5), Cm(0.5),
                    f"  ▸  {t}", font_size=13, color=C_DARK_TEXT)

    # Step 5A
    add_rect(slide, Cm(1.2), Cm(6.5), Cm(15.5), Cm(5.5), C_LIGHT_BG)
    add_rect(slide, Cm(1.2), Cm(6.5), Cm(1.5), Cm(5.5), C_BLUE)
    add_textbox(slide, Cm(1.25), Cm(7.3), Cm(1.4), Cm(0.7),
                "5A", font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3.1), Cm(6.6), Cm(13), Cm(0.7),
                "Step 5A ｜ 技术路线空间识别", font_size=14, bold=True, color=C_NAVY)
    items5a = ["识别可行技术路径集合，而非直接推荐方案",
               "列出候选路线及其约束条件",
               "分析各路线的风险与运维成本",
               "不对外发布此阶段内容"]
    for j, t in enumerate(items5a):
        add_textbox(slide, Cm(3.3), Cm(7.4 + j * 0.65), Cm(13), Cm(0.6),
                    f"  ▸  {t}", font_size=12, color=C_DARK_TEXT)

    # Step 5B
    add_rect(slide, Cm(17.5), Cm(6.5), Cm(15.7), Cm(5.5), C_LIGHT_BG)
    add_rect(slide, Cm(17.5), Cm(6.5), Cm(1.5), Cm(5.5), C_NAVY)
    add_textbox(slide, Cm(17.55), Cm(7.3), Cm(1.4), Cm(0.7),
                "5B", font_size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(19.4), Cm(6.6), Cm(13.5), Cm(0.7),
                "Step 5B ｜ 技术路线裁决（决策核心）", font_size=14, bold=True, color=C_NAVY)
    add_textbox(slide, Cm(19.6), Cm(7.4), Cm(13.2), Cm(0.55),
                "裁决依据（允许）", font_size=12, bold=True, color=C_OK)
    for j, t in enumerate(["风险可控性", "运维能力匹配度", "投资合理性"]):
        add_textbox(slide, Cm(19.8), Cm(7.95 + j * 0.55), Cm(12.8), Cm(0.5),
                    f"  ✔  {t}", font_size=12, color=C_OK)
    add_textbox(slide, Cm(19.6), Cm(9.65), Cm(13.2), Cm(0.55),
                "禁止依据", font_size=12, bold=True, color=C_WARN)
    for j, t in enumerate(["技术先进性", "厂商推荐", "模型默认倾向"]):
        add_textbox(slide, Cm(19.8), Cm(10.2 + j * 0.55), Cm(12.8), Cm(0.5),
                    f"  ✕  {t}", font_size=12, color=C_WARN)

    # 底部
    add_rect(slide, Cm(1.2), Cm(12.3), Cm(32.0), Cm(1.5), C_NAVY)
    add_textbox(slide, Cm(1.8), Cm(12.45), Cm(30), Cm(0.6),
                "裁决结果写入 specs/technical_baseline.md 并冻结。",
                font_size=14, bold=True, color=C_WHITE)
    add_textbox(slide, Cm(1.8), Cm(13.05), Cm(30), Cm(0.55),
                "冻结后，任何方案内容不得偏离此基线。",
                font_size=13, color=RGBColor(0xAA, 0xBB, 0xCC))


def slide_10_steps678(prs):
    slide = blank_slide(prs)
    draw_header(slide, "验证阶段详解（Step 6–8）", "Validation Phase")
    draw_footer_line(slide)
    draw_page_num(slide, 10)

    blocks = [
        ("6", "Step 6 ｜ 原厂询价 / 厂商测算（现实验证）",
         ["目的：验证设备市场可获得性，验证造价区间合理性",
          "避免评审阶段出现价格质疑",
          "产出：厂商报价单（归档至 specs/ 或附件）"],
         C_OK),
        ("7", "Step 7 ｜ 投资模型与整体报价构建",
         ["形成 CAPEX 结构（软硬件 / 集成 / 服务费分类）",
          "设备与扩容模型，每项配置有对应需求依据",
          "预算可解释路径（每项有合同与清单依据）",
          "产出：specs/investment_schema.md（金额冻结）"],
         C_BLUE),
        ("8", "Step 8 ｜ 投资合理性评审（高风险阶段）",
         ["重点审查：是否存在过度设计（配置超出业务需求）",
          "是否符合财政投资逻辑",
          "是否易被专家否决（单价偏高 / 数量不合理）",
          "⚠ Step 8 通过后，投资模型不再调整"],
         C_WARN),
    ]
    for i, (num, title, items, col) in enumerate(blocks):
        y = Cm(2.9 + i * 4.3)
        add_rect(slide, Cm(1.2), y, Cm(32.0), Cm(4.0), C_LIGHT_BG)
        add_rect(slide, Cm(1.2), y, Cm(1.5), Cm(4.0), col)
        add_textbox(slide, Cm(1.25), y + Cm(1.1), Cm(1.4), Cm(0.85),
                    num, font_size=26, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, Cm(3.1), y + Cm(0.2), Cm(29), Cm(0.7),
                    title, font_size=14, bold=True, color=C_NAVY)
        for j, t in enumerate(items):
            c = C_WARN if t.startswith("⚠") else C_DARK_TEXT
            add_textbox(slide, Cm(3.3), y + Cm(1.0 + j * 0.7), Cm(28.5), Cm(0.65),
                        f"  ▸  {t}", font_size=13, color=c)


def slide_11_steps910(prs):
    slide = blank_slide(prs)
    draw_header(slide, "输出阶段详解（Step 9–10）", "Output Phase")
    draw_footer_line(slide)
    draw_page_num(slide, 11)

    # Step 9
    add_rect(slide, Cm(1.2), Cm(2.9), Cm(32.0), Cm(7.0), C_LIGHT_BG)
    add_rect(slide, Cm(1.2), Cm(2.9), Cm(1.5), Cm(7.0), C_NAVY)
    add_textbox(slide, Cm(1.25), Cm(4.8), Cm(1.4), Cm(0.7),
                "9", font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3.1), Cm(3.05), Cm(29), Cm(0.7),
                "Step 9 ｜ 建设方案书写（表达阶段）", font_size=15, bold=True, color=C_NAVY)
    add_textbox(slide, Cm(3.1), Cm(3.85), Cm(29), Cm(0.55),
                "进入条件（全部满足，方可开始）：", font_size=13, bold=True, color=C_BLUE)
    cond = ["specs/需求.md 已完成并通过评审，包含现网设备清单",
            "specs/technical_baseline.md 已冻结",
            "specs/investment_schema.md 已确认无冲突"]
    for j, t in enumerate(cond):
        add_textbox(slide, Cm(3.5), Cm(4.5 + j * 0.68), Cm(28.5), Cm(0.6),
                    f"  ☑  {t}", font_size=13, color=C_OK)
    add_textbox(slide, Cm(3.1), Cm(6.65), Cm(29), Cm(0.55),
                "写作行为约束：", font_size=13, bold=True, color=C_BLUE)
    constraints = ["不新增需求", "不扩展未裁决能力", "不改变既定决策逻辑"]
    for j, t in enumerate(constraints):
        add_textbox(slide, Cm(3.5), Cm(7.25 + j * 0.62), Cm(28.5), Cm(0.55),
                    f"  ✕  {t}", font_size=13, color=C_WARN)
    add_rect(slide, Cm(3.1), Cm(9.0), Cm(29.5), Cm(0.75), C_NAVY)
    add_textbox(slide, Cm(3.3), Cm(9.05), Cm(29), Cm(0.6),
                "方案文本用于表达决策，而非探索决策。",
                font_size=13, bold=True, color=C_WHITE)

    # Step 10
    add_rect(slide, Cm(1.2), Cm(10.3), Cm(32.0), Cm(4.0), C_LIGHT_BG)
    add_rect(slide, Cm(1.2), Cm(10.3), Cm(1.5), Cm(4.0), C_ACCENT)
    add_textbox(slide, Cm(1.25), Cm(11.6), Cm(1.4), Cm(0.7),
                "10", font_size=22, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(3.1), Cm(10.45), Cm(29), Cm(0.7),
                "Step 10 ｜ 综合评审与一致性复核（终审）", font_size=15, bold=True, color=C_NAVY)
    dims = [("技术逻辑一致性", "章节内容与 technical_baseline 一致"),
            ("投资逻辑一致性", "预算数字与 investment_schema 完全吻合"),
            ("章节结构一致性", "各章节之间不存在相互矛盾"),
            ("可答辩性", "专家提问有据可查，有文件支撑")]
    for j, (k, v) in enumerate(dims):
        add_textbox(slide, Cm(3.3), Cm(11.15 + j * 0.72), Cm(12), Cm(0.65),
                    f"  ▸  {k}", font_size=13, bold=True, color=C_NAVY)
        add_textbox(slide, Cm(15.5), Cm(11.15 + j * 0.72), Cm(17), Cm(0.65),
                    v, font_size=13, color=C_GRAY_TEXT)

    add_textbox(slide, Cm(1.5), H - Cm(1.5), Cm(30), Cm(0.6),
                "终审通过后，方可进入交付与发布流程。",
                font_size=12, color=C_GRAY_TEXT)


def slide_12_freeze(prs):
    slide = blank_slide(prs)
    draw_header(slide, "技术路线冻结机制", "Technical Route Freeze")
    draw_footer_line(slide)
    draw_page_num(slide, 12)

    # 冻结含义
    add_textbox(slide, Cm(1.2), Cm(3.0), Cm(32), Cm(0.6),
                "▌ 冻结的含义", font_size=14, bold=True, color=C_NAVY)
    add_textbox(slide, Cm(1.5), Cm(3.65), Cm(31), Cm(0.65),
                "技术路线一旦写入 specs/technical_baseline.md，即视为工程裁决结果：",
                font_size=13, color=C_DARK_TEXT)
    freeze_items = [
        "方案正文不得偏离基线文件",
        "图件与部署结构不得引入基线外节点",
        "章节描述不得进行替代路线比较",
        "不得在方案中引入演进性技术变体或架构分支",
    ]
    for j, t in enumerate(freeze_items):
        add_textbox(slide, Cm(2.0), Cm(4.35 + j * 0.65), Cm(30), Cm(0.6),
                    f"  ✕  {t}", font_size=13, color=C_WARN)

    # 基线表格
    add_textbox(slide, Cm(1.2), Cm(7.0), Cm(32), Cm(0.6),
                "▌ 当前技术路线基线（示例）", font_size=14, bold=True, color=C_NAVY)
    rows = [
        ["层级", "裁决结论"],
        ["计算架构", "国产化 x86 架构服务器体系"],
        ["操作系统", "国产服务器操作系统"],
        ["数据库", "国产数据库管理系统（达梦等）"],
        ["存储模型", "NAS 统一文件存储资源池"],
        ["设备层定位", "数据采集接入层，不承担计算能力"],
    ]
    draw_table(slide, Cm(1.2), Cm(7.7), Cm(22),
               [Cm(6), Cm(16)], rows, font_size=12, row_h=Cm(0.65))

    # 排除项
    add_textbox(slide, Cm(24.0), Cm(7.0), Cm(9.5), Cm(0.6),
                "▌ 明确排除的技术模型", font_size=13, bold=True, color=C_NAVY)
    excluded = [
        "分布式对象存储集群架构",
        "FC-SAN 集中式存储",
        "云存储 / 云资源池模型",
        "混合异构存储架构",
    ]
    for j, t in enumerate(excluded):
        add_rect(slide, Cm(24.0), Cm(7.7 + j * 1.1), Cm(9.5), Cm(0.95),
                 RGBColor(0xFD, 0xF2, 0xF2))
        add_textbox(slide, Cm(24.2), Cm(7.75 + j * 1.1), Cm(9.0), Cm(0.8),
                    f"  ✕  {t}", font_size=12, color=C_WARN)

    add_rect(slide, Cm(1.2), Cm(12.2), Cm(32.0), Cm(1.3), C_NAVY)
    add_textbox(slide, Cm(1.8), Cm(12.35), Cm(30), Cm(1.0),
                "以上排除项不属于\"未来可选\"，属于本项目既定范围外。方案中不得出现，图件中不得表达。",
                font_size=13, color=C_WHITE)


def slide_13_investment(prs):
    slide = blank_slide(prs)
    draw_header(slide, "投资控制逻辑", "Investment Control Logic")
    draw_footer_line(slide)
    draw_page_num(slide, 13)

    # 法律地位说明
    add_rect(slide, Cm(1.2), Cm(3.0), Cm(32.0), Cm(2.0), C_LIGHT_BG)
    add_rect(slide, Cm(1.2), Cm(3.0), Cm(0.2), Cm(2.0), C_ACCENT)
    add_textbox(slide, Cm(1.8), Cm(3.1), Cm(30), Cm(0.65),
                "specs/investment_schema.md 是本项目资源配置与预算的唯一合法来源",
                font_size=14, bold=True, color=C_NAVY)
    add_textbox(slide, Cm(1.8), Cm(3.75), Cm(30), Cm(0.65),
                "包括：服务器数量 / 设备配置参数 / 存储容量规模 / 预算设备清单 / 部署结构",
                font_size=13, color=C_GRAY_TEXT)

    # 禁止行为表格
    add_textbox(slide, Cm(1.2), Cm(5.3), Cm(32), Cm(0.6),
                "▌ 禁止行为", font_size=14, bold=True, color=C_NAVY)
    rows = [
        ["禁止行为", "风险说明"],
        ["自行计算服务器数量", "与清单冲突，引发评审质疑"],
        ["自行修改配置规模", "破坏预算平衡，影响合同依据"],
        ["自行扩展硬件类别", "新增未批准采购项，违反财政纪律"],
        ["基于经验推导容量与性能规模", "无合同依据，无法答辩"],
        ["修改总计金额", "直接违反财政合规要求"],
    ]
    draw_table(slide, Cm(1.2), Cm(6.0), Cm(32.0),
               [Cm(14), Cm(18)], rows, font_size=12, row_h=Cm(0.68))

    # 计算约束
    add_textbox(slide, Cm(1.2), Cm(11.5), Cm(32), Cm(0.6),
                "▌ 计算约束（不可违反）", font_size=14, bold=True, color=C_NAVY)
    add_rect(slide, Cm(1.2), Cm(12.1), Cm(32.0), Cm(2.3),
             RGBColor(0x1E, 0x1E, 0x2E))
    code = [
        ("各分项金额之和  =  总计", True, C_YELLOW),
        ("若计算冲突  →  以总计为最高约束", False, C_WHITE),
        ("不允许因模型推理调整预算结构", False, RGBColor(0xCC, 0xCC, 0xCC)),
        ("单价是政策合规参数，不是估算参数", False, RGBColor(0xCC, 0xCC, 0xCC)),
    ]
    for j, (t, bold, col) in enumerate(code):
        add_textbox(slide, Cm(2.0), Cm(12.25 + j * 0.52), Cm(30), Cm(0.48),
                    t, font_size=12, bold=bold, color=col, font="Courier New")


def slide_14_charts(prs):
    slide = blank_slide(prs)
    draw_header(slide, "图件体系说明", "Chart System Classification")
    draw_footer_line(slide)
    draw_page_num(slide, 14)

    # 定位说明
    add_rect(slide, Cm(1.2), Cm(2.9), Cm(32.0), Cm(1.2), C_NAVY)
    add_textbox(slide, Cm(1.8), Cm(3.05), Cm(30), Cm(0.9),
                "图件属于技术决策表达资产，不是自由内容生成资产。用于表达既定决策，解释系统结构，支撑评审理解。",
                font_size=13, color=C_WHITE)

    # 四类图件
    cats = [
        ("A 类（强制）", "架构决策类", "investment_schema / 需求.md",
         C_WARN, "A-01 系统总体架构图\nA-02 网络拓扑图\nA-03 数据流转图\nA-04 部署结构图\nA-05 服务器角色划分图\nA-06 现网结构图"),
        ("B 类", "实施与流程类", "已裁决流程",
         C_BLUE, "B-01 实施甘特图\nB-02 系统切换流程图\nB-03 数据迁移流程图\nB-04 风险控制闭环图"),
        ("C 类", "投标增强类", "仅投标文件启用",
         C_OK, "C-01 技术响应映射图\nC-02 服务保障体系图\nC-03 质量控制闭环图"),
        ("D 类", "UI 表达类", "output/screenshot/",
         C_GRAY_TEXT, "D-01 页面结构原型图\nD-02 页面布局结构图\nD-03 页面流程交互图\nD-04 系统界面示意图\nD-05 现网系统截图"),
    ]
    for i, (label, name, src, col, items) in enumerate(cats):
        x = Cm(1.2) + i * Cm(8.1)
        w = Cm(7.8)
        add_rect(slide, x, Cm(4.35), w, Cm(0.85), col)
        add_textbox(slide, x, Cm(4.38), w, Cm(0.75),
                    label, font_size=13, bold=True, color=C_WHITE,
                    align=PP_ALIGN.CENTER)
        add_rect(slide, x, Cm(5.2), w, Cm(7.0), C_LIGHT_BG)
        add_textbox(slide, x + Cm(0.2), Cm(5.3), w - Cm(0.3), Cm(0.65),
                    name, font_size=13, bold=True, color=C_NAVY)
        add_textbox(slide, x + Cm(0.2), Cm(5.95), w - Cm(0.3), Cm(0.55),
                    f"来源：{src}", font_size=11, color=C_GRAY_TEXT)
        add_textbox(slide, x + Cm(0.2), Cm(6.6), w - Cm(0.3), Cm(5.5),
                    items, font_size=11, color=C_DARK_TEXT)

    # A-06 特别说明
    add_rect(slide, Cm(1.2), Cm(12.65), Cm(32.0), Cm(1.45),
             RGBColor(0xFD, 0xF7, 0xE7))
    add_rect(slide, Cm(1.2), Cm(12.65), Cm(0.2), Cm(1.45), C_YELLOW)
    add_textbox(slide, Cm(1.7), Cm(12.75), Cm(30), Cm(0.55),
                "A-06 现网结构图特别说明：",
                font_size=12, bold=True, color=C_NAVY)
    add_textbox(slide, Cm(1.7), Cm(13.3), Cm(30), Cm(0.6),
                "事实来源唯一为 specs/需求.md，禁止引用 investment_schema.md 或正文描述推导。",
                font_size=12, color=C_DARK_TEXT)


def slide_15_chart_rules(prs):
    slide = blank_slide(prs)
    draw_header(slide, "图件生成纪律", "Chart Generation Discipline")
    draw_footer_line(slide)
    draw_page_num(slide, 15)

    # 格式纪律表
    add_textbox(slide, Cm(1.2), Cm(2.9), Cm(20), Cm(0.6),
                "▌ 格式纪律（不可违反）", font_size=14, bold=True, color=C_NAVY)
    rows = [
        ["格式", "用途", "是否允许"],
        ["HTML", "正式图件源格式（唯一合法）", "✔"],
        ["PNG / SVG", "派生交付格式（必须基于 HTML 导出）", "✔"],
        ["Mermaid", "推理草图（仅限 drafts/chart/）", "✔"],
        ["ASCII 图", "—", "✕ 禁止"],
        ["直接生成 PNG", "跳过 HTML", "✕ 禁止"],
    ]
    draw_table(slide, Cm(1.2), Cm(3.6), Cm(19.0),
               [Cm(5), Cm(10.5), Cm(3.5)], rows, font_size=12, row_h=Cm(0.65))
    add_textbox(slide, Cm(1.2), Cm(7.55), Cm(19), Cm(0.55),
                "PNG 分辨率要求：≥ 1920×1080", font_size=12,
                bold=True, color=C_ACCENT)

    # 右侧：执行顺序
    add_textbox(slide, Cm(21.0), Cm(2.9), Cm(12.5), Cm(0.6),
                "▌ 图件生成执行顺序（强制）", font_size=14, bold=True, color=C_NAVY)
    steps = [
        "① 完成正文结构生成",
        "② 根据 chart_registry.md 判断触发条件",
        "③ 生成 drafts/chart/ 推理图",
        "④ 校验一致性",
        "   现状类 ↔ specs/需求.md",
        "   建设类 ↔ investment_schema",
        "⑤ 生成 output/chart/ 正式 HTML 图件",
        "⑥ 导出 PNG",
        "⑦ 正文引用图件",
    ]
    for j, t in enumerate(steps):
        col = C_GRAY_TEXT if t.startswith("   ") else C_DARK_TEXT
        bold = t[0] in "①②③④⑤⑥⑦"
        add_textbox(slide, Cm(21.0), Cm(3.6 + j * 0.62), Cm(12.5), Cm(0.58),
                    t, font_size=13, bold=bold, color=col)

    add_rect(slide, Cm(21.0), Cm(9.25), Cm(12.5), Cm(0.75), C_WARN)
    add_textbox(slide, Cm(21.2), Cm(9.3), Cm(12.0), Cm(0.6),
                "正文未完成前，禁止生成正式图件",
                font_size=12, bold=True, color=C_WHITE)

    # 冲突裁决
    add_textbox(slide, Cm(1.2), Cm(8.1), Cm(32), Cm(0.6),
                "▌ 图件冲突裁决优先级", font_size=14, bold=True, color=C_NAVY)
    priority = [
        ("1", "architecture_scope", C_WARN),
        ("2", "technical_baseline", RGBColor(0xE0, 0x70, 0x30)),
        ("3", "investment_schema", C_BLUE),
        ("4", "chart_registry", C_ACCENT),
        ("5", "正文描述", C_GRAY_TEXT),
    ]
    for j, (num, name, col) in enumerate(priority):
        x = Cm(1.2) + j * Cm(6.3)
        add_rect(slide, x, Cm(8.8), Cm(6.0), Cm(1.5), col)
        add_textbox(slide, x, Cm(8.9), Cm(6.0), Cm(0.6),
                    f"  {num}",
                    font_size=18, bold=True, color=C_WHITE)
        add_textbox(slide, x, Cm(9.5), Cm(6.0), Cm(0.65),
                    f"  {name}", font_size=12, color=C_WHITE)
        if j < 4:
            add_textbox(slide, x + Cm(6.0), Cm(9.2), Cm(0.3), Cm(0.65),
                        ">", font_size=16, color=C_GRAY_LINE)


def slide_16_writing_rules(prs):
    slide = blank_slide(prs)
    draw_header(slide, "写作行为约束", "Writing Behavior Constraints")
    draw_footer_line(slide)
    draw_page_num(slide, 16)

    # 六条纪律
    add_textbox(slide, Cm(1.2), Cm(2.9), Cm(32), Cm(0.6),
                "▌ 六条强制写作纪律", font_size=14, bold=True, color=C_NAVY)
    rules = [
        "不得整体重写既有章节",
        "不得改变章节标题与编号结构",
        "若章节已存在，仅允许增量补充",
        "禁止删除用户未明确要求删除的内容",
        "禁止扩展用户未提出的系统能力",
        "禁止解释写作理由或模型决策过程",
    ]
    for j, r in enumerate(rules):
        x = Cm(1.2) + (j % 3) * Cm(10.9)
        y = Cm(3.65) + (j // 3) * Cm(2.2)
        add_rect(slide, x, y, Cm(10.5), Cm(1.9), C_LIGHT_BG)
        add_rect(slide, x, y, Cm(0.2), Cm(1.9), C_WARN)
        add_textbox(slide, x + Cm(0.5), y + Cm(0.15), Cm(9.8), Cm(0.65),
                    f"✕  规则 {j+1}", font_size=11, bold=True, color=C_WARN)
        add_textbox(slide, x + Cm(0.5), y + Cm(0.8), Cm(9.7), Cm(1.0),
                    r, font_size=13, color=C_DARK_TEXT)

    # 章节风格要求
    add_textbox(slide, Cm(1.2), Cm(8.4), Cm(32), Cm(0.6),
                "▌ 章节风格对照", font_size=14, bold=True, color=C_NAVY)
    rows = [
        ["允许", "禁止"],
        ["功能性章节", "论文式结构"],
        ["直接服务项目建设逻辑", "背景化散文写法"],
        ["客观中性表达", "营销 / 宣传 / 产品介绍风格"],
        ["结构清晰、可被摘录审查", "情绪化表达"],
        ["方案级抽象描述", "工程设计书风格 / 算法细节"],
    ]
    draw_table(slide, Cm(1.2), Cm(9.1), Cm(32.0),
               [Cm(16), Cm(16)], rows, font_size=13, row_h=Cm(0.65))

    add_textbox(slide, Cm(1.5), H - Cm(1.5), Cm(30), Cm(0.65),
                "信息密度原则：删除无功能语句，优先结构清晰而非修辞，每段必须承担明确文档功能。",
                font_size=12, color=C_GRAY_TEXT)


def slide_17_violations1(prs):
    slide = blank_slide(prs)
    draw_header(slide, "常见违规示例对比（上）", "Common Violations — Part 1")
    draw_footer_line(slide)
    draw_page_num(slide, 17)

    groups = [
        ("场景一：需求未固化，直接进入写作", [
            ("收到用户口述后，直接让 AI 生成建设方案", "先澄清需求，固化为 specs/需求.md，通过评审后再生成"),
            ("方案内容来自会议讨论，未形成结构化文件", "需求必须以文件形式固化，口头共识不能替代文档"),
        ]),
        ("场景二：技术路线未冻结，出现多方案对比", [
            ("方案正文中写\"可采用 A 方案或 B 方案\"", "技术路线已裁决，正文只写裁决结果"),
            ("章节中出现\"也可考虑对象存储\"", "对象存储已被明确排除，不得出现在方案中"),
            ("以\"技术更先进\"为由推荐非基线技术", "先进性不是裁决依据，合规性与可控性才是"),
        ]),
        ("场景三：投资数字与 investment_schema 不一致", [
            ("正文写\"共配置 3 台应用服务器\"（与清单不符）", "服务器数量必须来自 investment_schema，不得自行推导"),
            ("根据存储容量需求反推硬盘数量", "硬盘数量以 investment_schema 清单为准"),
            ("预算总计与各分项之和不等", "各分类小计之和必须严格等于总计，不允许误差"),
        ]),
    ]

    y = Cm(2.9)
    for group_title, pairs in groups:
        add_rect(slide, Cm(1.2), y, Cm(32.0), Cm(0.7), C_NAVY)
        add_textbox(slide, Cm(1.5), y + Cm(0.07), Cm(31), Cm(0.55),
                    group_title, font_size=13, bold=True, color=C_WHITE)
        y += Cm(0.75)
        for wrong, right in pairs:
            row_h = Cm(0.9)
            add_rect(slide, Cm(1.2), y, Cm(15.8), row_h,
                     RGBColor(0xFD, 0xF2, 0xF2))
            add_textbox(slide, Cm(1.5), y + Cm(0.1), Cm(15.3), row_h - Cm(0.1),
                        f"✕  {wrong}", font_size=12, color=C_WARN)
            add_rect(slide, Cm(17.4), y, Cm(15.8), row_h,
                     RGBColor(0xF0, 0xF7, 0xF0))
            add_textbox(slide, Cm(17.7), y + Cm(0.1), Cm(15.3), row_h - Cm(0.1),
                        f"✔  {right}", font_size=12, color=C_OK)
            y += row_h + Cm(0.08)
        y += Cm(0.25)


def slide_18_violations2(prs):
    slide = blank_slide(prs)
    draw_header(slide, "常见违规示例对比（下）", "Common Violations — Part 2")
    draw_footer_line(slide)
    draw_page_num(slide, 18)

    groups = [
        ("场景四：图件越权", [
            ("拓扑图中出现 investment_schema 未定义的服务器", "图件节点必须来自 investment_schema，不得新增"),
            ("现网结构图参考了建设规划内容", "现网结构图唯一来源是 specs/需求.md"),
            ("图件中出现完整 IP 地址", "拓扑图、部署图不得出现完整 IP 地址"),
            ("将 Mermaid 草图直接作为正式图件交付", "正式图件必须为 HTML 格式，PNG 为派生交付格式"),
        ]),
        ("场景五：写作行为越权", [
            ("补充章节时重写了原有段落", "既有章节只允许增量补充，不得整体重写"),
            ("方案中扩展了\"未来可扩展为 AI 分析平台\"", "仅表达已裁决能力，不引入未来扩展描述"),
            ("章节末尾加入'建议后续...'", "禁止建议性语气和未来能力设想"),
        ]),
        ("场景六：现网与建设层混用", [
            ("现状分析章节引用了建设后的设备规模", "现状描述仅来自 specs/需求.md"),
            ("建设方案以现网存在的问题推导配置规模", "建设配置来自 investment_schema，不以现网问题推导"),
        ]),
    ]

    y = Cm(2.9)
    for group_title, pairs in groups:
        add_rect(slide, Cm(1.2), y, Cm(32.0), Cm(0.7), C_NAVY)
        add_textbox(slide, Cm(1.5), y + Cm(0.07), Cm(31), Cm(0.55),
                    group_title, font_size=13, bold=True, color=C_WHITE)
        y += Cm(0.75)
        for wrong, right in pairs:
            row_h = Cm(0.9)
            add_rect(slide, Cm(1.2), y, Cm(15.8), row_h,
                     RGBColor(0xFD, 0xF2, 0xF2))
            add_textbox(slide, Cm(1.5), y + Cm(0.1), Cm(15.3), row_h - Cm(0.1),
                        f"✕  {wrong}", font_size=12, color=C_WARN)
            add_rect(slide, Cm(17.4), y, Cm(15.8), row_h,
                     RGBColor(0xF0, 0xF7, 0xF0))
            add_textbox(slide, Cm(17.7), y + Cm(0.1), Cm(15.3), row_h - Cm(0.1),
                        f"✔  {right}", font_size=12, color=C_OK)
            y += row_h + Cm(0.08)
        y += Cm(0.25)


def slide_19_checklist(prs):
    slide = blank_slide(prs)
    draw_header(slide, "提交前检查清单", "Pre-Submission Checklist")
    draw_footer_line(slide)
    draw_page_num(slide, 19)

    sections = [
        ("一、事实层（specs/ 完整性）", [
            "specs/需求.md 已完成，包含\"现网设备清单\"章节",
            "specs/technical_baseline.md 已冻结，状态为\"已裁决\"",
            "specs/investment_schema.md 已确认，总计 = 各分项之和",
            "specs/chart_registry.md 已确认本项目适用图件类别",
        ], C_BLUE),
        ("二、技术路线一致性", [
            "方案正文未出现基线排除的技术模型",
            "方案正文未进行技术路线对比或替代方案描述",
            "图件中无未定义节点，无完整 IP 地址",
        ], C_ACCENT),
        ("三、投资合理性", [
            "正文中设备数量与 investment_schema 一致",
            "预算章节金额与 investment_schema 完全一致",
            "无过度设计描述（配置超出业务需求的解释）",
        ], C_OK),
        ("四、图件体系", [
            "A 类强制图件已全部生成（A-01 至 A-06）",
            "正式图件均为 HTML 源格式，存放于 output/chart/",
            "PNG 均为 HTML 导出，分辨率 ≥ 1920×1080",
        ], C_NAVY),
        ("五、写作合规性", [
            "无整体重写既有章节行为",
            "无营销 / 宣传 / 论文式语言",
            "无未来能力设想或\"建议后续\"表达",
        ], C_WARN),
        ("六、可答辩性", [
            "每项技术选择有 specs/ 文件支撑",
            "每项预算有清单依据",
            "章节结构无内部矛盾",
        ], C_GRAY_TEXT),
    ]

    for i, (title, items, col) in enumerate(sections):
        x = Cm(1.0) + (i % 3) * Cm(10.9)
        y = Cm(2.9) + (i // 3) * Cm(5.0)
        w = Cm(10.6)
        h = Cm(4.7)
        add_rect(slide, x, y, w, h, C_LIGHT_BG)
        add_rect(slide, x, y, w, Cm(0.8), col)
        add_textbox(slide, x + Cm(0.2), y + Cm(0.1), w - Cm(0.3), Cm(0.6),
                    title, font_size=11, bold=True, color=C_WHITE)
        for j, item in enumerate(items):
            add_textbox(slide, x + Cm(0.2), y + Cm(0.95 + j * 0.88), w - Cm(0.3),
                        Cm(0.82),
                        f"  ☐  {item}", font_size=11, color=C_DARK_TEXT)

    add_rect(slide, Cm(1.0), H - Cm(1.6), Cm(32.0), Cm(0.85), C_WARN)
    add_textbox(slide, Cm(1.5), H - Cm(1.55), Cm(31), Cm(0.7),
                "以上任何一项不满足，禁止进入交付流程。",
                font_size=13, bold=True, color=C_WHITE)


def slide_20_summary(prs):
    slide = blank_slide(prs)
    draw_header(slide, "总结", "Summary")
    draw_footer_line(slide)
    draw_page_num(slide, 20)

    # 价值总结表
    add_textbox(slide, Cm(1.2), Cm(2.9), Cm(32), Cm(0.6),
                "▌ Proposal Writing Skill 的核心工程价值", font_size=14,
                bold=True, color=C_NAVY)
    rows = [
        ["控制维度", "解决的问题"],
        ["需求固化机制", "避免口头需求带来的理解偏差与后期扯皮"],
        ["技术路线冻结", "避免方案在写作过程中技术漂移"],
        ["投资模型权威", "避免评审阶段出现价格质疑与数字冲突"],
        ["图件注册控制", "避免图件越权表达未批准的系统能力"],
        ["写作行为约束", "避免模型自由发挥破坏方案内部一致性"],
    ]
    draw_table(slide, Cm(1.2), Cm(3.6), Cm(32.0),
               [Cm(9), Cm(23)], rows, font_size=13, row_h=Cm(0.72))

    # 三个禁止
    add_textbox(slide, Cm(1.2), Cm(9.0), Cm(32), Cm(0.6),
                "▌ 记住三个\"禁止\"", font_size=14, bold=True, color=C_NAVY)
    prohibits = [
        "禁止在需求未固化前进入写作",
        "禁止在技术路线未冻结前生成图件与正文",
        "禁止模型自行推导投资数字与设备配置",
    ]
    for j, p in enumerate(prohibits):
        add_rect(slide, Cm(1.2), Cm(9.7 + j * 1.0), Cm(32.0), Cm(0.9), C_WARN)
        add_textbox(slide, Cm(1.6), Cm(9.78 + j * 1.0), Cm(31), Cm(0.72),
                    f"  ✕  {p}", font_size=14, bold=True, color=C_WHITE)

    # 核心定义
    add_rect(slide, Cm(1.2), Cm(12.85), Cm(32.0), Cm(2.2), C_NAVY)
    add_textbox(slide, Cm(1.8), Cm(13.0), Cm(30), Cm(0.55),
                "▌ 记住一个定义", font_size=12, bold=True,
                color=RGBColor(0xAA, 0xBB, 0xCC))
    add_textbox(slide, Cm(1.8), Cm(13.55), Cm(30), Cm(0.72),
                "方案 = 决策与约束的表达载体",
                font_size=20, bold=True, color=C_YELLOW)
    add_textbox(slide, Cm(1.8), Cm(14.25), Cm(30), Cm(0.65),
                "写方案不是在创作，而是把\"已经确定的事情\"用符合评审语境的方式说清楚。",
                font_size=13, color=RGBColor(0xCC, 0xCC, 0xCC))

    add_textbox(slide, Cm(1.5), H - Cm(1.55), Cm(30), Cm(0.6),
                "内部培训材料 · 请勿对外发布                依据文件：SKILL.md / readme.md / specs/ 体系文件",
                font_size=11, color=C_GRAY_TEXT)


# ══════════════════════════════════════════════════════════
# 主程序
# ══════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_01_cover(prs)
    slide_02_overview(prs)
    slide_03_positioning(prs)
    slide_04_principles(prs)
    slide_05_file_structure(prs)
    slide_06_authority(prs)
    slide_07_tensteps_overview(prs)
    slide_08_steps123(prs)
    slide_09_steps45(prs)
    slide_10_steps678(prs)
    slide_11_steps910(prs)
    slide_12_freeze(prs)
    slide_13_investment(prs)
    slide_14_charts(prs)
    slide_15_chart_rules(prs)
    slide_16_writing_rules(prs)
    slide_17_violations1(prs)
    slide_18_violations2(prs)
    slide_19_checklist(prs)
    slide_20_summary(prs)

    out = "output/Proposal_Writing_Skill_培训PPT.pptx"
    prs.save(out)
    print("Done: " + out + " (" + str(len(prs.slides)) + " slides)")


if __name__ == "__main__":
    main()
